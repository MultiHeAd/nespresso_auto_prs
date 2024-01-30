# -*- coding:utf-8 -*-
from openpyxl import load_workbook
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
import datetime
import pandas as pd
import numpy as np
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
import os

def get_dt_year_month():
    dt = datetime.datetime.now()
    year, month = dt.year, dt.month - 1
    if month == 0:
        month = 12
        year -= 1
    return f"{year}.{str(month).zfill(2)}"

def get_dt_year_month_last():
    dt = datetime.datetime.now()
    year, month = dt.year, dt.month - 1
    if month == 0:
        month = 12
        year -= 2
    return f"{year}.{str(month).zfill(2)}"

def get_dt_yearNmonth():
    dt = datetime.datetime.now()
    year, month = dt.year, dt.month - 1
    if month == 0:
        month = 12
        year -= 1
    return year, month

def format_rpt_title(prs, i, j):
    shape = prs.slides[i].shapes[j]
    year, month = get_dt_yearNmonth()
    shape.text = f'Nespresso\n{year}年{month}月月报'
    print(f'format_rpt_title--{i}--{j}--{shape.text}')

def format_table_data(prs, fdata, slide_number, table_index, sheet_name, start_row, end_row, start_col, end_col):
    slide_number += 1
    skiprows = start_row - 1
    nrows = end_row - start_row + 1
    usecols = f"{start_col}:{end_col}"
    df = pd.read_excel(fdata, sheet_name=sheet_name, skiprows=skiprows, nrows=nrows, usecols=usecols)

    color_df = df.applymap(lambda x: RGBColor(73, 113, 30) if isinstance(x, (int, float)) and x > 0 else RGBColor(192,0,0))

    df = df.applymap(lambda x: f"{int(x * 100)}%" if isinstance(x, (int, float)) else x)
    #print("读取到的 DataFrame:\n", df)

    slide = prs.slides[slide_number - 1]
    table = slide.shapes[table_index].table

    for i, row in enumerate(df.itertuples(index=False)):
        for j, val in enumerate(row):
            cell = table.cell(i+1, j)
            cell.text = str(val)

            # 设置字体、字号、居中对齐和颜色
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
                for run in paragraph.runs:
                    run.font.name = '微软雅黑'
                    run.font.size = Pt(10.5)

                    if j == len(row) - 1:  # 最后一列
                        run.font.color.rgb = color_df.iloc[i, j]

    return prs
def format_chart_title(prs, i, j, lbl):
    chart = prs.slides[i].shapes[j].chart
    year, month = get_dt_yearNmonth()
    chart.chart_title.text_frame.paragraphs[0].text = f'{year}年{month}月{lbl}'
    print(f'format_chart_title--{i}--{j}')

def format_shape_text(prs, i, j, time_lbl, cont):
    shape = prs.slides[i].shapes[j]
    shape.text_frame.paragraphs[0].runs[0].text = f'{time_lbl}{cont}' if cont != '' else f'{time_lbl}'
    print(f'format_shape_text--{i}--{j}--{shape.text}')

def format_table_title(prs, i, j, cell_i, cell_j, time_lbl, cont):
    table = prs.slides[i].shapes[j].table
    cell = table.cell(cell_i, cell_j)
    cell.text_frame.paragraphs[0].runs[0].text = f'{time_lbl}{cont}'
    print(f'format_table_title--{i}--{j}--{cell.text}')

def get_tbl_data(fdata, table_name, col_c, col_s, nrows, ct=''):
    # df = pg.read_table(table_name)
    df = pd.read_excel(fdata, sheet_name=table_name)
    if ct != '':
        df = df[df['类目']==ct].reset_index(drop=True)
    else:
        df = df
    dfcategories = df[col_c][-nrows:].fillna('').reset_index(drop=True)
    dfseries = df[col_s][-nrows:].fillna(0).reset_index(drop=True)
    dfseries = dfseries.apply(lambda x: x*1)
    return dfcategories, dfseries

def get_tbl_data9(fdata, table_name, col_c, col_s, nrows, ct=''):
    # df = pg.read_table(table_name)
    df = pd.read_excel(fdata, sheet_name=table_name)
    if ct != '':
        df = df[df['类目']==ct].reset_index(drop=True)
    else:
        df = df
    dfcategories = df[col_c][-nrows:].fillna('').reset_index(drop=True)
    dfseries = df[col_s][-nrows:].fillna(0).reset_index(drop=True)
    dfseries = dfseries.apply(lambda x: x * 10)
    return dfcategories, dfseries
def get_tbl_data10(fdata, table_name, col_c, col_s, nrows, ct=''):
    # df = pg.read_table(table_name)
    df = pd.read_excel(fdata, sheet_name=table_name)
    if ct != '':
        df = df[df['类目']==ct].reset_index(drop=True)
    else:
        df = df
    dfcategories = df[col_c][-nrows:].fillna('').reset_index(drop=True)
    dfseries = df[col_s][-nrows:].fillna(0).reset_index(drop=True)
    dfseries = dfseries.apply(lambda x: x / 1000)
    return dfcategories, dfseries


def get_tbl_data_chl(fdata, table_name, col_c, col_s, chl='', ct='', shop=''):
    # df = pg.read_table(table_name)
    df = pd.read_excel(fdata, sheet_name=table_name)
    if ct != '':
        df = df[df['类目']==ct].reset_index(drop=True)
    else:
        df = df
    if shop != '':
        df = df[df['店铺']==shop].reset_index(drop=True)
    else:
        df = df
    if chl == '广告流量':
        df = df[df['一级']==chl].reset_index(drop=True)
        df = df[df['三级']!='汇总'].sort_values(by=['rk_ad']).reset_index(drop=True)
        dfcategories = df[col_c].fillna('').reset_index(drop=True)
        if '二级' in dfcategories.columns:
            dfcategories['二级'].mask(dfcategories['二级'].shift(1) == dfcategories['二级'], inplace=True)
            dfcategories = dfcategories.fillna('')
        else:
            dfcategories = dfcategories
        dfseries = df[col_s].fillna(0).reset_index(drop=True)
    elif chl == '平台流量':
        df = df[df['一级']==chl].reset_index(drop=True)
        df = df[df['三级']=='汇总'].reset_index(drop=True)
        df = df[df['二级']!='汇总'].sort_values(by=['rk_uv']).reset_index(drop=True)
        dfcategories = df[col_c][:7].fillna('').reset_index(drop=True)
        dfseries = df[col_s][:7].fillna(0).reset_index(drop=True)
    return dfcategories, dfseries

def get_tbl_data_chl_1st(fdata, table_name, col_c, col_s, chl='', ct='', shop=''):
    # df = pg.read_table(table_name)
    df = pd.read_excel(fdata, sheet_name=table_name)
    if ct != '':
        df = df[df['类目']==ct].reset_index(drop=True)
    else:
        df = df
    if shop != '':
        df = df[df['店铺'].isin(shop)].reset_index(drop=True)
        df['店铺'] = df['店铺'].astype('category').cat.reorder_categories(shop)
        # df['店铺'].cat.reorder_categories(shop, inplace=True)
        df.sort_values('店铺', inplace=True)
    else:
        df = df
    if chl == '广告流量':
        df = df[df['一级']==chl].reset_index(drop=True)
        df = df[df['二级']=='汇总'].reset_index(drop=True)
        df = df[df['三级']=='汇总'].sort_values(by=['rk_ad']).reset_index(drop=True)
        dfcategories = df[col_c]
        dfseries = df[col_s].fillna(0).reset_index(drop=True)
    elif chl == '平台流量':
        df = df[df['一级']==chl].reset_index(drop=True)
        df = df[df['二级']=='汇总'].reset_index(drop=True)
        df = df[df['三级']=='汇总'].sort_values(by=['rk_uv']).reset_index(drop=True)
        dfcategories = df[col_c]
        dfseries = df[col_s].fillna(0).reset_index(drop=True)
    return dfcategories, dfseries

def get_tbl_data_chl_sort(fdata, table_name, col_c, col_s, chl='', ct='', shop=''):
    df = pd.read_excel(fdata, sheet_name=table_name)
    if ct != '':
        df = df[df['类目']==ct].reset_index(drop=True)
    else:
        df = df
    if shop != '':
        df = df[df['店铺']==shop].reset_index(drop=True)
    else:
        df = df
    if chl == '广告流量':
        df = df[df['一级']==chl].reset_index(drop=True)
        df = df[df['三级']!='汇总'].sort_values(by=['rk_ad']).reset_index(drop=True)

        df=df.groupby('二级',sort=False,as_index=False).apply(lambda x: x.sort_values('访客数_本月', ascending=False)).reset_index(drop=True)

        dfcategories = df[col_c].fillna('').reset_index(drop=True)
        if '二级' in dfcategories.columns:
            dfcategories['二级'].mask(dfcategories['二级'].shift(1) == dfcategories['二级'], inplace=True)
            dfcategories = dfcategories.fillna('')
        else:
            dfcategories = dfcategories
        dfseries = df[col_s].fillna(0).reset_index(drop=True)
    elif chl == '平台流量':
        df = df[df['一级']==chl].reset_index(drop=True)
        df = df[df['三级']=='汇总'].reset_index(drop=True)
        df = df[df['二级']!='汇总'].sort_values(by=['rk_uv']).reset_index(drop=True)
        dfcategories = df[col_c][:7].fillna('').reset_index(drop=True)
        dfseries = df[col_s][:7].fillna(0).reset_index(drop=True)
    return dfcategories, dfseries

def replace_chart_data(chart, dfcategories, dfseries):
    chart_data = CategoryChartData()
    if len(dfcategories.keys()) > 1:
        cols = dfcategories.keys()
        if '标签' in cols:
            dfcategories['标签'].mask(dfcategories['标签'].shift(1) == dfcategories['标签'], inplace=True)
            dfcategories = dfcategories.fillna('')
        else:
            dfcategories = dfcategories.fillna('')
        if '二级' in cols:
            dfcategories['二级'].mask(dfcategories['二级'].shift(1) == dfcategories['二级'], inplace=True)
            dfcategories = dfcategories.fillna('')
        else:
            dfcategories = dfcategories.fillna('')
        for ct in dfcategories.to_dict(orient='records'):
            chart_data.add_category(f"{ct[cols[0]]}").add_sub_category(f"{ct[cols[1]]}")
    else:
        chart_data.categories = dfcategories.iloc[:, 0]
    for col_s in dfseries.keys():
        col_s = str(col_s)
        if '_本月' in col_s or '-本月' in col_s or '当月' in col_s:
            col_s_fin = f"{get_dt_year_month()[2:]}{col_s.replace('_本月', '', ).replace('-本月', '', ).replace('当月', '', )}"
        elif '同比月' in col_s:
            col_s_fin = f"{get_dt_year_month_last()[2:]}{col_s.replace('_同比月', '', ).replace('-同比月', '', ).replace('同比月', '', )}"
        else:
            col_s_fin = col_s
        chart_data.add_series(col_s_fin, dfseries[col_s])
    chart.replace_data(chart_data)

def format_chart_data(prs, fdata, i, j, table_name, col_c, col_s, nrows, ct=''):
    chart = prs.slides[i].shapes[j].chart
    dfcategories, dfseries = get_tbl_data(fdata, table_name, col_c, col_s, nrows, ct)
    replace_chart_data(chart, dfcategories, dfseries)
    print(f'format_chart_data--{i}--{j}')

def format_chart_data9(prs, fdata, i, j, table_name, col_c, col_s, nrows, ct=''):
    chart = prs.slides[i].shapes[j].chart
    dfcategories, dfseries = get_tbl_data9(fdata, table_name, col_c, col_s, nrows, ct)
    replace_chart_data(chart, dfcategories, dfseries)
    print(f'format_chart_data--{i}--{j}')
def format_chart_data10(prs, fdata, i, j, table_name, col_c, col_s, nrows, ct=''):
    chart = prs.slides[i].shapes[j].chart
    dfcategories, dfseries = get_tbl_data10(fdata, table_name, col_c, col_s, nrows, ct)
    replace_chart_data(chart, dfcategories, dfseries)
    print(f'format_chart_data--{i}--{j}')

def format_chart_line(prs, i, j, series_num):
    chart = prs.slides[i].shapes[j].chart
    chart.series[series_num].format.line.width = 12700
    chart.series[series_num].format.line.color.theme_color = 15
    chart.series[series_num].format.line.dash_style = 2
    print(f'format_chart_line--{i}--{j}')

def format_chart_data2(prs, fdata, i, j, table_name, col_c, col_s, chl='', ct='', shop=''):
    chart = prs.slides[i].shapes[j].chart
    dfcategories, dfseries = get_tbl_data_chl(fdata, table_name, col_c, col_s, chl, ct, shop)
    replace_chart_data(chart, dfcategories, dfseries)
    print(f'format_chart_data2--{i}--{j}')

def format_chart_data_sort(prs, fdata, i, j, table_name, col_c, col_s, chl='', ct='', shop=''):
    chart = prs.slides[i].shapes[j].chart
    dfcategories, dfseries = get_tbl_data_chl_sort(fdata, table_name, col_c, col_s, chl, ct, shop)
    replace_chart_data(chart, dfcategories, dfseries)
    print(f'format_chart_data2--{i}--{j}')

def format_chart_data_chl(prs, fdata, i, j, table_name, col_c, col_s, chl='', ct='', shop=''):
    chart = prs.slides[i].shapes[j].chart
    dfcategories1, dfseries1 = get_tbl_data_chl_1st(fdata, table_name, col_c, col_s, chl[0], ct, shop)
    dfcategories2, dfseries2 = get_tbl_data_chl_1st(fdata, table_name, col_c, col_s, chl[1], ct, shop)
    dfcategories = pd.DataFrame()
    dfcategories[col_c] = dfcategories1[col_c].astype('O').replace("旗舰店|电器旗舰店|官方旗舰店|家享咖啡|delonghi", "", regex=True).replace("barsetto", "百胜图", regex=True).replace("Dolce Gusto", "多趣酷思", regex=True)
    dfseries = pd.concat([dfseries1, dfseries2], axis=1)
    dfseries.columns = chl
    replace_chart_data(chart, dfcategories, dfseries)
    print(f'format_chart_data_chl--{i}--{j}')

def get_tbl_data_trans(fdata, table_name, col, col_s, nrows, ct=''):
    # df = pg.read_table(table_name)
    df = pd.read_excel(fdata, sheet_name=table_name)
    if ct != '':
        df = df[df['类目'] == ct].reset_index(drop=True)
    else:
        df = df
    df_trans = df[col_s][-nrows:].T.reset_index() if nrows == 1 else df[col_s][-nrows:1-nrows].T.reset_index()
    df_trans.columns = ['tag', col]
    dfcategories = df_trans[['tag']].replace(f"{col}", "", regex=True).replace("占比", "", regex=True).replace("top品牌", "", regex=True).replace("品牌", "", regex=True).replace("_", "", regex=True).replace("nd", "", regex=True).replace("rd", "", regex=True)
    dfcategories['tag']=dfcategories['tag'].str.upper()
    dfseries = df_trans[[col]].fillna(0)
    return dfcategories, dfseries


def format_chart_data3(prs, fdata, i, j, table_name, col, col_s, nrows, ct=''):
    chart = prs.slides[i].shapes[j].chart
    dfcategories, dfseries = get_tbl_data_trans(fdata, table_name, col, col_s, nrows, ct)
    replace_chart_data(chart, dfcategories, dfseries)
    print(f'format_chart_data3--{i}--{j}')

def get_tbl_data2(fdata, table_name, col_c, col_s, col_l, nrows, ct=''):
    # df = pg.read_table(table_name)
    df = pd.read_excel(fdata, sheet_name=table_name)
    if ct != '':
        df = df[df['类目']==ct].reset_index(drop=True)
    else:
        df = df
    dfcategories = df[col_c][-nrows:].fillna('').reset_index(drop=True) if nrows==1 else df[col_c][-nrows:1-nrows].fillna('').reset_index(drop=True)
    dfseries = df[col_s][-nrows:].fillna(0).reset_index(drop=True) if nrows==1 else df[col_s][-nrows:1-nrows].fillna(0).reset_index(drop=True)
    dflables = df[col_l][-nrows:].fillna(0).reset_index(drop=True) if nrows==1 else df[col_l][-nrows:1-nrows].fillna(0).reset_index(drop=True)
    return dfcategories, dfseries, dflables

def replace_chart_data_label(chart, dfseries, dflabels, col_ex):
    plot = chart.plots[0]
    for i, series in enumerate(plot.series):
        if i in col_ex:continue
        for j, point in enumerate(series.points):
            data_label = point.data_label
            data_label.text_frame.text = f'{"{:.0%}".format(dflabels.iloc[j][i])}\n{"{:,.0f}".format(dfseries.iloc[j][i])}'

def format_chart_data4(prs, fdata, i, j, table_name, col_c, col_s, col_l, col_ex, nrows, ct=''):
    chart = prs.slides[i].shapes[j].chart
    dfcategories1, dfseries1, dflabels1 = get_tbl_data2(fdata, table_name, col_c, col_s, col_l, nrows[0], ct)
    dfcategories2, dfseries2, dflabels2 = get_tbl_data2(fdata, table_name, col_c, col_s, col_l, nrows[1], ct)
    dfcategories = pd.concat([dfcategories1,dfcategories2])
    dfseries = pd.concat([dfseries1, dfseries2])
    dflabels = pd.concat([dflabels1, dflabels2])
    dfseries.columns = [str(x).split('_', 1)[-1] for x in dfseries.columns]
    replace_chart_data(chart, dfcategories, dfseries)
    replace_chart_data_label(chart, dfseries, dflabels, col_ex)
    print(f'format_chart_data4--{i}--{j}')

def format_chart_data5(prs, fdata, i, j, table_name, col, col_s, nrows, ct=''):
    chart = prs.slides[i].shapes[j].chart
    dfcategories1, dfseries1 = get_tbl_data_trans(fdata, table_name, col, col_s, nrows[0], ct)
    dfcategories2, dfseries2 = get_tbl_data_trans(fdata, table_name, col, col_s, nrows[1], ct)

    dfseries1.index = dfseries2.index

    dfseries = pd.concat([dfseries1, dfseries2], axis=1)
    dfseries.columns = [get_dt_year_month_last(), get_dt_year_month()]

    dfcategories = pd.DataFrame()
    dfcategories[col] = dfcategories1['tag'].replace("_", "", regex=True).replace("", '店铺整体')

    replace_chart_data(chart, dfcategories, dfseries)
    print(f'format_chart_data5--{i}--{j}')

def get_tbl_data3(fdata, table_name, col_c, col_s, col_l, nrows, ct=''):
    # df = pg.read_table(table_name)
    df = pd.read_excel(fdata, sheet_name=table_name)
    if ct != '':
        df = df[df['类目'] == ct].reset_index(drop=True)
    else:
        df = df
    dfcategories = df[col_c][:nrows].fillna('').reset_index(drop=True)
    dfseries = df[col_s][:nrows].fillna(0).reset_index(drop=True)
    dflabels = df[col_l][:nrows].fillna('').reset_index(drop=True)
    return dfcategories, dfseries, dflabels

def replace_chart_data_label2(chart, dflabels):
    plot = chart.plots[0]
    for i, series in enumerate(plot.series):
        for j, point in enumerate(series.points):
            data_label = point.data_label
            data_label.text_frame.text = dflabels.iloc[j][i]

def format_chart_data6(prs, fdata, i, j, table_name, col_c, col_s, col_l, nrows, ct=''):
    chart = prs.slides[i].shapes[j].chart
    dfcategories, dfseries, dflabels = get_tbl_data3(fdata, table_name, col_c, col_s, col_l, nrows, ct)
    replace_chart_data(chart, dfcategories, dfseries)
    replace_chart_data_label2(chart, dflabels)
    print(f'format_chart_data6--{i}--{j}')

def data_val_format(data_val, data_key):
    if data_val == 0:
        data_val_fin = '持平' if data_key=='排名_变化' else '-'
    elif data_key=='排名_变化':
        data_val_fin= "{:+.0f}".format(data_val)
    elif data_key in ['*首图']:
        data_val_fin = ""
    elif data_val in ['<2000', '-'] or data_key in ['一级类目购买偏好', 'TGI', '叶子类目Top1', '叶子类目Top1 (按TGI)', '偏好购买品牌Top 5（按TGI排序）', '购买率(本品vs.竞品)', '本品竞争力', '购买率\n(本品vs.竞品)', 'lbl', 'TOP 单品（SPU）', '咖啡机类型', 'TOP5单品', '型号系列', 'TOP10单品', '品牌', '搜索词']:
        data_val_fin = str(data_val)
    else:
        if data_key in ['时间范围']:
            data_val_fin = data_val.replace('-', '.')
        elif data_key in ['chl']:
            data_val_fin = data_val.replace('品销宝- ', '')
        elif data_key in ['支付转化率_同比']:
            data_val_fin = "{:+.0f}pp".format(data_val * 100) if abs(data_val) >= 0.02 else "{:+.1f}pp".format(data_val * 100)
        elif data_key in ['品牌市占率', '支付转化率', '增速'] or '占比' in data_key:
            if data_val > 9.99:
                data_val_fin = ">999%"
            elif data_val < -9.99:
                data_val_fin = "<-999%"
            else:
                data_val_fin = "{:.0%}".format(data_val) if abs(data_val) >= 0.02 else "{:.1%}".format(data_val)
        elif '同比' in data_key:
            if data_val > 9.99:
                data_val_fin = ">999%"
            elif data_val < -9.99:
                data_val_fin = "<-999%"
            else:
                data_val_fin = "{:.0%}".format(data_val) if abs(data_val) >= 0.02 else "{:.1%}".format(data_val)
        else:
            data_val_fin = "{:,.0f}".format(data_val)
    return data_val_fin

def format_table_cell(cell, data_key, data_val_fin, ct):
    if not cell.text_frame.paragraphs:
        cell.text_frame.add_paragraph()
    paragraph = cell.text_frame.paragraphs[0]

    if not paragraph.runs:
        paragraph.add_run()
    run = paragraph.runs[0]

    run.text = data_val_fin
    if data_val_fin in ['-', '持平']:
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
        print(f'format_table_cell--{cell.text}--{cell.text_frame.paragraphs[0].runs[0].font.color.rgb}')
    elif data_key not in['品牌GMV同比','官旗GMV同比','非官旗GMV同比'] and data_key in ['排名_变化', '支付转化率'] or '同比' in data_key:
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(192, 0, 0) if '-' in data_val_fin else RGBColor(73, 113, 30)
        print(f'format_table_cell--{cell.text}--{cell.text_frame.paragraphs[0].runs[0].font.color.rgb}')
    elif data_key in['品牌GMV同比','官旗GMV同比','非官旗GMV同比']:
        if '-' in data_val_fin:
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(192, 0, 0)
        else:
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(73, 113, 30)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(219, 240, 197)
        print(f'format_table_cell--{cell.text}--{cell.text_frame.paragraphs[0].runs[0].font.color.rgb}')
    elif data_key in ['uv_本月']:
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(89, 89, 89)
        print(f'format_table_cell--{cell.text}--{cell.text_frame.paragraphs[0].runs[0].font.color.rgb}')
    elif data_key in ['uv']:
        cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(165, 165, 165)
        print(f'format_table_cell--{cell.text}--{cell.text_frame.paragraphs[0].runs[0].font.color.rgb}')
    else:
        if ct == '咖啡机':
            if '同比' in data_key:
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(192, 0, 0) if '-' in data_val_fin else RGBColor(0, 0, 0)
            else:
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
            print(f'format_table_cell--{cell.text}--{cell.text_frame.paragraphs[0].runs[0].font.color.rgb}')
        elif ct == '胶囊咖啡':
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(0, 0, 0)
            if '同比' in data_key and '-' not in data_val_fin:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(219, 240, 197)
                print(f'format_table_cell--{cell.text}--{cell.text_frame.paragraphs[0].runs[0].font.color.rgb}--{cell.text_frame.paragraphs[0].runs[0].font.color.rgb}')




def format_table_data_profile(prs, fdata, i, j, table_name, col_c, col_s, nrows, ct):
    table = prs.slides[i].shapes[j].table
    dfcategories, dfseries = get_tbl_data(fdata, table_name, col_c, col_s, nrows, ct)
    for ii, row in enumerate(table.rows):
        for jj, cell in enumerate(row.cells):
            if jj<=1:continue
            data_key = dfseries.iloc[:, ii].name
            data_val = dfseries.iloc[jj - 2, ii]
            data_val_fin = data_val_format(data_val, data_key)
            format_table_cell(cell, data_key, data_val_fin, ct)
    print(f'format_table_data_profile--{i}--{j}')

def format_table_data_chl(prs, fdata, i, j, table_name, col_c, col_s, chl, ct='', shop=''):
    table = prs.slides[i].shapes[j].table
    dfcategories, dfseries = get_tbl_data_chl(fdata, table_name, col_c, col_s, chl, ct, shop)
    dfcategories_1st, dfseries_1st = get_tbl_data_chl_1st(fdata, table_name, col_c, col_s, chl, ct)
    for ii, row in enumerate(table.rows):
        for jj, cell in enumerate(row.cells):
            if jj<=1:continue
            elif jj == 2:
                data_key = dfseries_1st.iloc[:, ii].name
                data_val = dfseries_1st.iloc[jj - 3, ii]
            else:
                data_key = dfseries.iloc[:, ii].name
                data_val = dfseries.iloc[jj - 3, ii]
            data_val_fin = data_val_format(data_val, data_key)
            #print(data_key, data_val, data_val_fin)
            format_table_cell(cell, data_key, data_val_fin, ct)
            
    print(f'format_table_data_chl--{i}--{j}')

def format_table_data_chl_sort(prs, fdata, i, j, table_name, col_c, col_s, chl, ct='', shop=''):
    table = prs.slides[i].shapes[j].table
    dfcategories, dfseries = get_tbl_data_chl_sort(fdata, table_name, col_c, col_s, chl, ct, shop)
    dfcategories_1st, dfseries_1st = get_tbl_data_chl_1st(fdata, table_name, col_c, col_s, chl, ct)
    #print(dfcategories,dfseries)
    col_no_pay=dfcategories['chl'][dfseries['支付人数_同比']==0].map(lambda x:x.split('-')[-1]).tolist()
    for ii, row in enumerate(table.rows):
        for jj, cell in enumerate(row.cells):
            if jj<=1:continue
            elif jj == 2:
                data_key = dfseries_1st.iloc[:, ii].name
                data_val = dfseries_1st.iloc[jj - 3, ii]
            else:
                data_key = dfseries.iloc[:, ii].name
                data_val = dfseries.iloc[jj - 3, ii]
            data_val_fin = data_val_format(data_val, data_key)
            #print(data_key, data_val, data_val_fin)
            format_table_cell(cell, data_key, data_val_fin, ct)
    print(f'format_table_data_chl--{i}--{j}')
    return col_no_pay

def format_table_data_brand(prs, fdata, i, j, table_name, col_c, col_s, nrows, ct):
    table = prs.slides[i].shapes[j].table
    dfcategories, dfseries, dflabels = get_tbl_data3(fdata, table_name, col_c, col_s, col_c, nrows, ct)
    for ii, row in enumerate(table.rows):
        for jj, cell in enumerate(row.cells):
            if ii<=1:continue
            data_key = dfseries.columns[jj]
            data_val = dfseries.iloc[ii - 2, jj]
            data_val_fin = '' if data_val == '' else data_val_format(data_val, data_key)
            format_table_cell(cell, data_key, data_val_fin, ct)
    print(f'format_table_data_brand--{i}--{j}')

def format_table_data_brand_rk(prs, fdata, i, j, table_name, col_c, col_s, nrows, ct):
    table = prs.slides[i].shapes[j].table
    dfcategories, dfseries, dflabels = get_tbl_data3(fdata, table_name, col_c, col_s, col_c, nrows, ct)
    for ii, row in enumerate(table.rows):
        for jj, cell in enumerate(row.cells):
            data_key = dfseries.columns[jj]
            data_val = dfseries.iloc[ii , jj]
            data_val_fin = data_val_format(data_val, data_key)
            format_table_cell(cell, data_key, data_val_fin, ct)
    print(f'format_table_data_brand_rk--{i}--{j}')

def format_table_data_brand_mkt_share(prs, fdata, i, j, table_name, col_c, col_s, nrows, ct):
    table = prs.slides[i].shapes[j].table
    dfcategories1, dfseries1, dflabels1 = get_tbl_data2(fdata, table_name, col_c, col_s, col_c, nrows[0], ct)
    dfcategories2, dfseries2, dflabels2 = get_tbl_data2(fdata, table_name, col_c, col_s, col_c, nrows[1], ct)
    dfseries = pd.concat([dfseries1,dfseries2])
    for ii, row in enumerate(table.rows):
        for jj, cell in enumerate(row.cells):
            if ii==0:continue
            data_key = dfseries.columns[jj]
            data_val = dfseries.iloc[ii -1, jj]
            data_val_fin = data_val_format(data_val, data_key)
            format_table_cell(cell, data_key, data_val_fin, ct)
    print(f'format_table_data_brand_mkt_share--{i}--{j}')

def format_table_data_brand_mkt(prs, fdata, i, j, table_name, col_c, col_s, nrows, ct):
    table = prs.slides[i].shapes[j].table
    dfcategories, dfseries = get_tbl_data(fdata, table_name, col_c, col_s, nrows, ct)
    for ii, row in enumerate(table.rows):
        for jj, cell in enumerate(row.cells):
            if ii!=1 or jj==1 or jj==2:continue
            if jj == 0:
                data_key = dfseries.columns[jj]
                data_val = dfseries.iloc[ii - 1, jj]
            elif jj >= 3:
                data_key = dfseries.columns[jj - 2]
                data_val = dfseries.iloc[ii - 1, jj - 2]
            data_val_fin = data_val_format(data_val, data_key)
            format_table_cell(cell, data_key, data_val_fin, ct)
    print(f'format_table_data_brand_mkt--{i}--{j}')

def format_table_data_ct_chl(prs, fdata, i, j, table_name, col_c, col_s, chl='', ct='', shop=''):
    table = prs.slides[i].shapes[j].table
    dfcategories, dfseries = get_tbl_data_chl(fdata, table_name, col_c, col_s, chl, ct, shop)
    for ii, row in enumerate(table.rows):
        for jj, cell in enumerate(row.cells):
            data_key = dfseries.columns[jj]
            data_val = dfseries.iloc[ii, jj]
            data_val_fin = data_val_format(data_val, data_key)
            format_table_cell(cell, data_key, data_val_fin, ct)
    print(f'format_table_data_ct_chl--{i}--{j}')

def format_table_data_ct_chl_ad(prs, fdata, i, j, table_name, col_c, col_s, chl='', ct='', shops=''):
    table = prs.slides[i].shapes[j].table
    dfcategories = pd.DataFrame()
    dfseries = pd.DataFrame()
    for shop in shops:
        dfc, dfs = get_tbl_data_chl(fdata, table_name, col_c,col_s, chl, ct, shop)
        dfcategories = dfc.fillna('')
        dfseries = pd.concat([dfseries, dfs], axis=1)
    for ii, row in enumerate(table.rows):
        for jj, cell in enumerate(row.cells):
            if ii<=3:continue
            if jj<=1:
                data_key = 'chl'
                data_val = dfcategories.iloc[ii - 4, jj]
                if cell.text_frame.text == '':continue
            else:
                data_key = dfseries.columns[jj - 2]
                data_val = dfseries.iloc[ii - 4, jj - 2]
            if cell.text_frame.text == '':continue
            data_val_fin = data_val_format(data_val, data_key)
            cell.text_frame.paragraphs[0].runs[0].text = data_val_fin
    print(f'format_table_data_ct_chl_ad--{i}--{j}')

def format_table_data_val(prs, fdata, i, j, tbl_i, tbl_j, table_name, col, nrows, ct='', shop='', chl=''):
    table = prs.slides[i].shapes[j].table
    cell = table.rows[tbl_i].cells[tbl_j]
    data_val = get_tbl_data_val(fdata, table_name, col, nrows, ct, shop, chl)
    data_val_fin = data_val_format(data_val, col)
    format_table_cell(cell, col, data_val_fin, ct)
    print(f'format_table_data_val--{i}--{j}')

def format_remark_text(prs, i, j, lbl, asterisk=''):
    shape = prs.slides[i].shapes[j]
    year_month = get_dt_year_month()
    if asterisk == '':
        shape.text = f'数据来源：{lbl}\n时间维度：{year_month}'
    elif i == 8:
        shape.text = f'数据来源：{lbl}\n时间维度：{year_month}\n*{asterisk}\n*相关系数R：取值范围在[-1,1]，-表示负相关，+表示正相关，绝对值越接近1说明变量之间的线性关系越强；一般来说，|R|＜0.3相关关系较弱可视为不相关，0.3<|R|<0.5为低度相关，0.5<|R|<0.8为中度相关，|R|>0.8为高度相关;'
    else:
        shape.text = f'数据来源：{lbl}\n时间维度：{year_month}\n*{asterisk}'
    print(f'format_remark_text--{i}--{j}')

def format_text_text(prs, i, j, lbl):
    shape = prs.slides[i].shapes[j]
    year, month = get_dt_yearNmonth()
    shape.text_frame.paragraphs[0].runs[0].text = f'{year}年{month}月 {lbl}'
    print(f'format_text_text--{i}--{j}--{shape.text}')

def format_arrow(prs, i, j, dat_i, dat_j):
    shape = prs.slides[i].shapes[j]
    dat = prs.slides[dat_i].shapes[dat_j].text
    shape.rotation = 180.0 if '-' in dat else 0.0
    shape.fill.fore_color.rgb = RGBColor(237, 173, 86) if '-' in dat else RGBColor(73, 113, 30)
    print(f'format_arrow--{i}--{j}--{dat}--{shape.rotation}--{shape.fill.fore_color.rgb}')

def format_arrow_reverse(prs, i, j, dat_i, dat_j):
    shape = prs.slides[i].shapes[j]
    dat = prs.slides[dat_i].shapes[dat_j].text
    shape.rotation = 0.0 if '-' in dat else 180.0
    shape.fill.fore_color.rgb = RGBColor(73, 113, 30) if '-' in dat else RGBColor(237, 173, 86)
    print(i, j, dat, shape.rotation, shape.fill.fore_color.rgb)
    print(f'format_arrow_reverse--{i}--{j}--{dat}--{shape.rotation}--{shape.fill.fore_color.rgb}')

def format_text_data_color(prs, i, j):
    shape = prs.slides[i].shapes[j]
    shape.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(192, 0, 0) if '-' in shape.text else RGBColor(73, 113, 30)
    print(f'format_text_data_color--{i}--{j}--{shape.text}--{shape.text_frame.paragraphs[0].runs[0].font.color.rgb}')

def get_tbl_data_val(fdata, table_name, col, nrows=1, ct='', shop='', chl=''):
    # df = pg.read_table(table_name)
    df = pd.read_excel(fdata, sheet_name=table_name)
    if ct != '':
        df = df[df['类目']==ct].reset_index(drop=True)
    else:
        df = df
    if shop != '':
        df = df[df['店铺']==shop].reset_index(drop=True)
    else:
        df = df
    if chl != '':
        df = df[df['chl']==chl].reset_index(drop=True)
    else:
        df = df
    data_val = list(df[col])[-nrows]
    return data_val

def get_tbl_data_val2(fdata, table_name, col, nrows, ct='', shop='', chl=''):
    dat_val_last = get_tbl_data_val(fdata, table_name, col, nrows[0], ct, shop, chl)
    dat_val = get_tbl_data_val(fdata, table_name, col, nrows[1], ct, shop, chl)
    dat_val_diff = dat_val-dat_val_last
    return dat_val_diff

def format_text_data(prs, fdata, i, j, table_name, col, nrows, ct=''):
    shape = prs.slides[i].shapes[j]
    data_val = get_tbl_data_val(fdata, table_name, col, nrows, ct) if type(nrows) != list else get_tbl_data_val2(fdata, table_name, col, nrows, ct)
    # if ct in ['咖啡机', '咖啡豆/粉'] and col in ['销售金额']:
    #     data_val_fin = "{:,.2f}亿".format(data_val/100000000)
    if col in ['购买频次','单次购买量']:
        data_val_fin = "{:,.2f}".format(data_val)
    # elif ct == '胶囊咖啡' and col in ['销售金额']:
    #     data_val_fin = "{:,.0f}万".format(data_val/10000)
    # elif col in ['销售金额_速溶咖啡', '销售金额_咖啡豆']:
    #     data_val_fin = "{:,.2f}".format(data_val/100000000)
    elif col in ['TTL*']:
        if data_val>100000000:
            data_val_fin = "{:,.2f}亿".format(data_val/100000000)
        else:
            data_val_fin = "{:,.0f}万".format(data_val / 10000)
    # elif col in ['销售金额_挂耳咖啡', '销售金额_咖啡液', '销售金额_胶囊咖啡', '销售金额_咖啡粉']:
    #     data_val_fin = "{:,.0f}".format(data_val/10000)
    elif col in ['支付转化率']:
        data_val_fin = "{:.0%}".format(data_val)
    elif col in ['支付转化率_同比', '老客gmv占比增幅', '老客人数占比增幅', '支付金额_胶囊咖啡_占比', '支付买家数_胶囊咖啡_占比']:
        data_val_fin = "{:+.1f}pp".format(data_val*100)
    elif '同比' in col:
        data_val_fin = "{:+.0%}".format(data_val) if abs(data_val) > 0.01 else "{:+.1%}".format(data_val)
    else:
        data_val_fin = "{:,.0f}".format(data_val)
    shape.text_frame.paragraphs[0].runs[0].text = data_val_fin
    print(f'format_text_data--{i}--{j}--{shape.text}')

def get_tbl_data_wb(fdata, sheet_name, usecols, skiprows, nrows, sep=1):
    df = pd.read_excel(fdata, sheet_name=sheet_name, usecols=usecols, skiprows=skiprows, nrows=nrows)
    df.columns = np.array(df).tolist()[0]
    df.drop([0], inplace=True)
    df.columns = [str(x) for x in df.columns]
    dfcategories = df.iloc[:, :sep].fillna('').reset_index(drop=True)
    dfseries = df.iloc[:, sep:].fillna(0).reset_index(drop=True)
    return dfcategories, dfseries

def get_tbl_data_wb2(fdata, sheet_name, usecols, skiprows, nrows, sep=1):
    df = pd.read_excel(fdata, sheet_name=sheet_name, usecols=usecols, skiprows=skiprows, nrows=nrows)

    new_header = df.iloc[0]  # 第一行作为标题
    df = df[1:]  # 取数据除去标题行
    df.columns = new_header  # 设置新标题

    dfcategories = df.iloc[:, :sep].reset_index(drop=True)
    dfseries = df.iloc[:, sep:].reset_index(drop=True)

    for col in dfseries.columns:
        dfseries[col] = pd.to_numeric(dfseries[col], errors='coerce').fillna(0)

    return dfcategories, dfseries

def format_chart_data_wb(prs, fdata, i, j, sheet_name, usecols, skiprows, nrows, sep=1):
    chart = prs.slides[i].shapes[j].chart
    dfcategories, dfseries = get_tbl_data_wb(fdata, sheet_name, usecols, skiprows, nrows, sep)
    replace_chart_data(chart, dfcategories, dfseries)
    print(f'format_chart_data_wb--{i}--{j}')

def format_chart_data_wb2(prs, fdata, i, j, sheet_name, usecols, skiprows, nrows, sep=1):
    chart = prs.slides[i].shapes[j].chart
    dfcategories, dfseries = get_tbl_data_wb2(fdata, sheet_name, usecols, skiprows, nrows, sep)
    replace_chart_data(chart, dfcategories, dfseries)
    print(f'format_chart_data_wb--{i}--{j}')

def format_chart_data_wb_item(prs, fdata, i, j, sheet_name, usecols, skiprows, nrows, sep=1):
    chart = prs.slides[i].shapes[j].chart
    dfcategories, dfseries = get_tbl_data_wb(fdata, sheet_name, usecols, skiprows, nrows, sep)
    dfnew = dfseries[dfseries[dfseries.columns[1]]>0].tail(14)
    if '月份' in dfnew.columns:
        dfnew['月份'] = dfnew['月份'].apply(lambda x: "{:.2f}".format(float(x)))
    dfcategories = dfnew.iloc[:, :sep]
    dfseries = dfnew.iloc[:, sep:]
    replace_chart_data(chart, dfcategories, dfseries)
    print(f'format_chart_data_wb_item--{i}--{j}')

def format_table_data_wb(prs, fdata, i, j, sheet_name, usecols, skiprows, nrows, skipi, skipj):
    table = prs.slides[i].shapes[j].table
    dfcategories, dfseries = get_tbl_data_wb(fdata, sheet_name, usecols, skiprows, nrows)
    for ii, row in enumerate(table.rows):
        for jj, cell in enumerate(row.cells):
            # print(ii,jj,cell.text_frame.text)
            if ii<=skipi:continue
            if jj<=skipj:continue
            # print(ii, jj, cell.text_frame.text)
            # print(ii, jj, cell.text_frame.text, "--", dfseries.iloc[ii - 1 - skipi, jj - skipj])
            data_key = dfseries.iloc[:, jj-skipj].name if str(dfseries.iloc[:, jj-skipj].name) != 'nan' else 'lbl'
            data_val = dfseries.iloc[ii - 1 - skipi, jj - skipj]
            # print(data_key, data_val)
            data_val_fin = data_val_format(data_val, data_key)
            # print(data_key, data_val, data_val_fin)
            format_table_cell(cell, data_key, data_val_fin, None)
    print(f'format_table_data_wb--{i}--{j}')

def format_table_cell_data_val(prs, i, j, cell_i, cell_j, pat):
    table = prs.slides[i].shapes[j].table
    cell = table.cell(cell_i, cell_j)
    cell.text_frame.paragraphs[0].runs[0].text = pat.format(float(cell.text_frame.text))
    print(f'format_table_cell_data_val--{i}--{j}--{cell.text}')

def format_text_data_wb(prs, fdata, i, j, sheet_name, usecols, skiprows, nrows, lbl):
    shape = prs.slides[i].shapes[j]
    dfcategories, dfseries = get_tbl_data_wb(fdata, sheet_name, usecols, skiprows, nrows)
    data_val = dfseries.iloc[0, 0]
    if lbl == 'TOP5GMV':
        data_val_fin = lbl + "：{:,.0f}w".format(int(data_val)/10000)
    elif lbl in ['Nespresso GMV占比']:
        data_val_fin = lbl + "：{:.0%}".format(float(data_val))
    elif lbl in ['cnt']:
        data_val_fin = "{:,.0f}".format(int(data_val))
    elif lbl in ['ratio']:
        data_val_fin = r"（{:.0%}）".format(float(data_val))
    shape.text_frame.paragraphs[0].runs[0].text = f"{data_val_fin}"
    print(f'format_text_data_wb--{i}--{j}--{shape.text}')

def format_table_cell_color(prs, i, j, cell_is, cell_js, lbl):
    table = prs.slides[i].shapes[j].table
    for cell_i in cell_is:
        for cell_j in cell_js:
            cell = table.cell(cell_i, cell_j)
            data_val = cell.text_frame.text
            if data_val != '-' and lbl == '品牌分销' and '+' in data_val:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(219, 240, 197)
                print(f'format_table_cell_color--{i}--{j}--{cell.fill.fore_color.rgb}')
            elif data_val != '-' and lbl == '本品竞争力':
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(73, 113, 30) if '高于竞品' in data_val else RGBColor(192, 0, 0) if '低于竞品' in data_val else RGBColor(0, 0, 0)
                print(f'format_table_cell_color--{i}--{j}--{cell.text_frame.paragraphs[0].runs[0].font.color.rgb}')
            elif data_val != '-' and lbl == 'TOP单品':
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 231, 231) if '-' in data_val else RGBColor(219, 240, 197)
                print(f'format_table_cell_color--{i}--{j}--{cell.fill.fore_color.rgb}')
            elif lbl == '咖啡机类型':
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = RGBColor(141, 121, 99) if '半自动' in data_val else RGBColor(115, 70, 31) if '全自动' in data_val else RGBColor(213, 140, 45)
                print(f'format_table_cell_color--{i}--{j}--{cell.text_frame.paragraphs[0].runs[0].font.color.rgb}')

def format_arrow_wb(prs, fdata, i, j, sheet_name, usecols, skiprows, nrows, sep=1):
    shape = prs.slides[i].shapes[j]
    dfcategories, dfseries = get_tbl_data_wb(fdata, sheet_name, usecols, skiprows, nrows, sep)
    dat = dfseries.iloc[0, 0]
    try:
        dat = float(dat)
    except ValueError:
        return
    shape.rotation = 180.0 if dat >= 0 else 0.0
    shape.fill.fore_color.rgb = RGBColor(73, 113, 30) if dat >= 0 else RGBColor(237, 173, 86)
    print(f'format_arrow_wb--{i}--{j}--{dat}--{shape.rotation}--{shape.fill.fore_color.rgb}')

def rps_format(prs, fdata):

    format_rpt_title(prs, 0, 0)

    format_text_text(prs, 3, 9, '销售金额')
    format_chart_data10(prs, fdata, 3, 6, 'dat_nespresso_shop_overview', ['时间范围'], ['支付金额'], 13)
    format_text_data(prs, fdata, 3, 16, 'dat_nespresso_shop_overview', '支付金额', 1)
    format_text_data(prs, fdata, 3, 17, 'dat_nespresso_shop_overview', '支付买家数', 1)
    format_text_data(prs, fdata, 3, 18, 'dat_nespresso_shop_overview', '访客数', 1)
    format_text_data(prs, fdata, 3, 19, 'dat_nespresso_shop_overview', '支付转化率', 1)
    format_text_data(prs, fdata, 3, 20, 'dat_nespresso_shop_overview', '客单价', 1)
    format_text_data(prs, fdata, 3, 29, 'dat_nespresso_shop_overview', '支付金额_同比', 1)
    format_text_data(prs, fdata, 3, 30, 'dat_nespresso_shop_overview', '支付买家数_同比', 1)
    format_text_data(prs, fdata, 3, 31, 'dat_nespresso_shop_overview', '访客数_同比', 1)
    format_text_data(prs, fdata, 3, 32, 'dat_nespresso_shop_overview', '支付转化率_同比', 1)
    format_text_data(prs, fdata, 3, 33, 'dat_nespresso_shop_overview', '客单价_同比', 1)
    format_text_data_color(prs, 3, 29)
    format_text_data_color(prs, 3, 30)
    format_text_data_color(prs, 3, 31)
    format_text_data_color(prs, 3, 32)
    format_text_data_color(prs, 3, 33)
    format_arrow(prs, 3, 15, 3, 29)
    format_arrow(prs, 3, 22, 3, 30)
    format_arrow(prs, 3, 24, 3, 31)
    format_arrow(prs, 3, 26, 3, 32)
    format_arrow(prs, 3, 28, 3, 33)
    format_remark_text(prs, 3, 2, '生意参谋')

    format_chart_data3(prs, fdata, 4, 3, 'dat_nespresso_shop_overview', 'gmv', ['新客gmv', '老客gmv'], 1)
    format_chart_data3(prs, fdata, 4, 4, 'dat_nespresso_shop_overview', '人数', ['新客人数', '老客人数'], 1)
    format_text_data(prs, fdata, 4, 7, 'dat_nespresso_shop_overview', '老客gmv占比增幅', 1)
    format_text_data(prs, fdata, 4, 10, 'dat_nespresso_shop_overview', '老客人数占比增幅', 1)
    format_text_data(prs, fdata, 4, 18, 'dat_nespresso_shop_overview', '新客gmv', 1)
    format_text_data(prs, fdata, 4, 20, 'dat_nespresso_shop_overview', '新客人数', 1)
    format_text_data(prs, fdata, 4, 22, 'dat_nespresso_shop_overview', '新客客单价', 1)
    format_text_data(prs, fdata, 4, 25, 'dat_nespresso_shop_overview', '新客gmv_同比', 1)
    format_text_data(prs, fdata, 4, 28, 'dat_nespresso_shop_overview', '新客人数_同比', 1)
    format_text_data(prs, fdata, 4, 31, 'dat_nespresso_shop_overview', '新客客单价_同比', 1)
    format_text_data(prs, fdata, 4, 33, 'dat_nespresso_shop_overview', '老客gmv', 1)
    format_text_data(prs, fdata, 4, 35, 'dat_nespresso_shop_overview', '老客人数', 1)
    format_text_data(prs, fdata, 4, 37, 'dat_nespresso_shop_overview', '老客客单价', 1)
    format_text_data(prs, fdata, 4, 40, 'dat_nespresso_shop_overview', '老客gmv_同比', 1)
    format_text_data(prs, fdata, 4, 43, 'dat_nespresso_shop_overview', '老客人数_同比', 1)
    format_text_data(prs, fdata, 4, 46, 'dat_nespresso_shop_overview', '老客客单价_同比', 1)
    format_text_data_color(prs, 4, 7)
    format_text_data_color(prs, 4, 10)
    format_text_data_color(prs, 4, 25)
    format_text_data_color(prs, 4, 28)
    format_text_data_color(prs, 4, 31)
    format_text_data_color(prs, 4, 40)
    format_text_data_color(prs, 4, 46)
    format_text_data_color(prs, 4, 43)
    format_arrow(prs, 4, 6, 4, 7)
    format_arrow(prs, 4, 9, 4, 10)
    format_arrow(prs, 4, 24, 4, 25)
    format_arrow(prs, 4, 27, 4, 28)
    format_arrow(prs, 4, 30, 4, 31)
    format_arrow(prs, 4, 39, 4, 40)
    format_arrow(prs, 4, 42, 4, 43)
    format_arrow(prs, 4, 45, 4, 46)
    format_remark_text(prs, 4, 2, '生意参谋', '新客定义：当月产生店铺购买行为但在前365天未发生店铺购买行为')

    format_chart_data4(prs, fdata, 5, 3, 'dat_nespresso_shop_overview', ['时间范围'], ['支付金额_咖啡机', '支付金额_胶囊咖啡', '支付金额_其他咖啡', '支付金额_其他'], ['支付金额_咖啡机_占比', '支付金额_胶囊咖啡_占比', '支付金额_其他咖啡_占比',  '支付金额_其他_占比'], [2], [13, 1])
    format_chart_data4(prs, fdata, 5, 4, 'dat_nespresso_shop_overview', ['时间范围'], ['支付买家数_咖啡机', '支付买家数_胶囊咖啡', '支付买家数_其他咖啡'], ['支付买家数_咖啡机_占比', '支付买家数_胶囊咖啡_占比', '支付买家数_其他咖啡_占比'], [2], [13, 1])
    format_chart_data5(prs, fdata, 5, 5, 'dat_nespresso_shop_overview', '客单价', ['客单价', '客单价_咖啡机',  '客单价_胶囊咖啡'], [13, 1])
    format_text_data(prs, fdata, 5, 8, 'dat_nespresso_shop_overview', '支付金额_胶囊咖啡_占比', [13, 1])
    format_text_data(prs, fdata, 5, 11, 'dat_nespresso_shop_overview', '支付买家数_胶囊咖啡_占比', [13, 1])
    format_text_data_color(prs, 5, 8)
    format_text_data_color(prs, 5, 11)
    format_arrow(prs, 5, 7, 5, 8)
    format_arrow(prs, 5, 10, 5, 11)
    format_arrow_reverse(prs, 5, 12, 5, 8)
    format_arrow_reverse(prs, 5, 13, 5, 11)
    format_table_data(prs, fdata, 5, 14, 'calculation', 1, 3, 'A', 'C')
    format_table_data(prs, fdata, 5, 15, 'calculation', 1, 3, 'E', 'G')
    format_remark_text(prs, 5, 2, '生意参谋')

    format_chart_data_sort(prs, fdata, 6, 4, 'dat_nespresso_shop_chl', ['二级', '三级'], ['访客数_同比月', '访客数_本月', '支付转化率_同比月', '支付转化率_本月'], '广告流量', '')
    col_no_pay=format_table_data_chl_sort(prs, fdata, 6, 3, 'dat_nespresso_shop_chl', ['chl'], ['支付人数_同比', '访客数_同比', '支付转化率_同比'], '广告流量', '')
    format_remark_text(prs, 6, 2, '生意参谋', '生意参谋暂不支持查看{}渠道的后链路转化表现'.format('、'.join(col_no_pay)))

    format_chart_data2(prs, fdata, 7, 3, 'dat_nespresso_shop_chl', ['二级'], ['访客数_同比月', '访客数_本月', '支付转化率_同比月', '支付转化率_本月'], '平台流量', '')
    format_table_data_chl(prs, fdata, 7, 4, 'dat_nespresso_shop_chl', ['chl'], ['支付人数_同比', '访客数_同比', '支付转化率_同比'], '平台流量', '')
    format_remark_text(prs, 7, 2, '生意参谋')

    format_chart_data_wb(prs, fdata, 8, 4, '店铺整体生意', 'C:E', 100, 11)
    format_chart_data_wb(prs, fdata, 8, 3, '店铺整体生意', 'C:E', 113, 11)
    format_remark_text(prs, 8, 2, '站内搜索词：生意参谋；小红书搜索词：千瓜数据', '小红书热度值中已剔除出行(酒店、民宿)及海外购机优惠(加拿大、Canada、羊毛)等相关笔记')

    format_chart_data9(prs, fdata, 11, 16, 'dat_nespresso_mkt_index', ['时间范围'], ['销售金额'], 14, '咖啡机')
    format_text_text(prs, 11, 8, '销售金额')
    format_text_data(prs, fdata, 11, 17, 'dat_nespresso_mkt_index', '销售金额', 1, '咖啡机')
    format_text_data(prs, fdata, 11, 19, 'dat_nespresso_mkt_index', '销量', 1, '咖啡机')
    format_text_data(prs, fdata, 11, 20, 'dat_nespresso_mkt_index', '购买人数', 1, '咖啡机')
    format_text_data(prs, fdata, 11, 21, 'dat_nespresso_mkt_index', '购买频次', 1, '咖啡机')
    format_text_data(prs, fdata, 11, 37, 'dat_nespresso_mkt_index', '单次购买量', 1, '咖啡机')
    format_text_data(prs, fdata, 11, 15, 'dat_nespresso_mkt_index', '平均价格', 1, '咖啡机')
    format_text_data(prs, fdata, 11, 18, 'dat_nespresso_mkt_index', '销售金额_同比', 1, '咖啡机')
    format_text_data(prs, fdata, 11, 24, 'dat_nespresso_mkt_index', '销量_同比', 1, '咖啡机')
    format_text_data(prs, fdata, 11, 27, 'dat_nespresso_mkt_index', '购买人数_同比', 1, '咖啡机')
    format_text_data(prs, fdata, 11, 30, 'dat_nespresso_mkt_index', '购买频次_同比', 1, '咖啡机')
    format_text_data(prs, fdata, 11, 39, 'dat_nespresso_mkt_index', '单次购买量_同比', 1, '咖啡机')
    format_text_data(prs, fdata, 11, 33, 'dat_nespresso_mkt_index', '平均价格_同比', 1, '咖啡机')
    format_text_data_color(prs, 11, 18)
    format_text_data_color(prs, 11, 24)
    format_text_data_color(prs, 11, 27)
    format_text_data_color(prs, 11, 30)
    format_text_data_color(prs, 11, 33)
    format_text_data_color(prs, 11, 39)
    format_arrow(prs, 11, 13, 11, 18)
    format_arrow(prs, 11, 23, 11, 24)
    format_arrow(prs, 11, 26, 11, 27)
    format_arrow(prs, 11, 29, 11, 30)
    format_arrow(prs, 11, 32, 11, 33)
    format_arrow(prs, 11, 36, 11, 39)
    format_remark_text(prs, 11, 2, '生意参谋；渠道：全渠道', '平均价格(指数)：商品的平均成交价格(为准确洞察，不包含0.01元的订单数据)；人均购买量(指数)：消费者平均一次购买的商品件数（=销量(指数)/购买人数(指数)/购买频次）；购买频次：每个消费者平均购买次数')

    format_chart_data3(prs, fdata, 12, 3, 'dat_nespresso_mkt_share', '销售金额', ['销售金额_top1品牌_占比', '销售金额_top2nd品牌_占比', '销售金额_top3rd品牌_占比', '销售金额_TOP4-10品牌_占比', '销售金额_TOP11-20品牌_占比', '销售金额_top品牌_其他商家_占比'], 1,'咖啡机')
    format_chart_data6(prs, fdata, 12, 14, 'dat_nespresso_mkt_rk_brand', ['品牌名称'], ['交易金额'], ['交易金额_万'], 10, '咖啡机')
    format_table_data_brand_mkt_share(prs, fdata,12, 4, 'dat_nespresso_mkt_share', ['时间范围'], ['时间范围', '销售金额_top3品牌_占比', '销售金额_top10品牌_占比'], [1, 13], '咖啡机')
    format_table_data_brand(prs, fdata,12, 7, 'dat_nespresso_mkt_rk_brand', ['品牌名称'], ['交易金额_同比', '品牌市占率', '品牌市占率_同比', '购买人数', '购买人数_同比', '客单价', '客单价_同比', '访客人数', '访客人数_同比'], 10, '咖啡机')
    #format_table_data_brand_mkt(prs, fdata,12, 7, 'dat_nespresso_mkt_overview', ['时间范围'], ['销售金额_同比', '购买人数', '购买人数_同比', '客单价', '客单价_同比', '访客人数', '访客人数_同比'], 1, '咖啡机')
    format_table_data_brand_rk(prs, fdata,12, 15, 'dat_nespresso_mkt_rk_brand', ['品牌名称'], ['排名_变化'], 10, '咖啡机')
    format_text_data(prs, fdata, 12, 9, 'dat_nespresso_mkt_rk_brand', 'TTL*', 1, '咖啡机')
    format_remark_text(prs, 12, 2, '生意参谋；渠道：全渠道','市场TTL：当月市场Top50品牌GMV加总；*市占率 = 品牌GMV/ 当月Top50品牌GMV加总')

    format_chart_data_wb(prs, fdata, 13, 4, '咖啡机表现', 'C:E', 69, 11)
    format_table_data_wb(prs, fdata, 13, 5, '咖啡机表现', [0, 1, 2, 5, 6, 7], 69, 11, 1, -1)
    format_table_cell_color(prs, 13, 5, list(range(2, 12)), list(range(1, 4)), '品牌分销')
    format_remark_text(prs, 13, 3, '生意参谋')

    format_chart_data(prs, fdata, 14, 6, 'dat_nespresso_shop_ct_each', ['时间范围'], ['支付金额'], 13, '咖啡机')
    format_text_text(prs, 14, 9, '销售金额')
    format_text_data(prs, fdata, 14, 16, 'dat_nespresso_shop_ct_each', '支付金额', 1, '咖啡机')
    format_text_data(prs, fdata, 14, 17, 'dat_nespresso_shop_ct_each', '支付买家数', 1, '咖啡机')
    format_text_data(prs, fdata, 14, 18, 'dat_nespresso_shop_ct_each', '商品访客数', 1, '咖啡机')
    format_text_data(prs, fdata, 14, 19, 'dat_nespresso_shop_ct_each', '支付转化率', 1, '咖啡机')
    format_text_data(prs, fdata, 14, 20, 'dat_nespresso_shop_ct_each', '客单价', 1, '咖啡机')
    format_text_data(prs, fdata, 14, 29, 'dat_nespresso_shop_ct_each', '支付金额_同比', 1, '咖啡机')
    format_text_data(prs, fdata, 14, 30, 'dat_nespresso_shop_ct_each', '支付买家数_同比', 1, '咖啡机')
    format_text_data(prs, fdata, 14, 31, 'dat_nespresso_shop_ct_each', '商品访客数_同比', 1, '咖啡机')
    format_text_data(prs, fdata, 14, 32, 'dat_nespresso_shop_ct_each', '支付转化率_同比', 1, '咖啡机')
    format_text_data(prs, fdata, 14, 33, 'dat_nespresso_shop_ct_each', '客单价_同比', 1, '咖啡机')
    format_text_data(prs, fdata, 14, 34, 'dat_nespresso_mkt_index', '销售金额_同比', 1, '咖啡机')
    format_text_data(prs, fdata, 14, 35, 'dat_nespresso_mkt_index', '购买人数_同比', 1, '咖啡机')

    format_text_data_color(prs, 14, 29)
    format_text_data_color(prs, 14, 30)
    format_text_data_color(prs, 14, 31)
    format_text_data_color(prs, 14, 32)
    format_text_data_color(prs, 14, 33)
    format_text_data_color(prs, 14, 34)
    format_text_data_color(prs, 14, 35)
    format_arrow(prs, 14, 15, 14, 29)
    format_arrow(prs, 14, 22, 14, 30)
    format_arrow(prs, 14, 24, 14, 31)
    format_arrow(prs, 14, 26, 14, 32)
    format_arrow(prs, 14, 28, 14, 33)
    format_remark_text(prs, 14, 2, '生意参谋')

    format_chart_data(prs, fdata, 16, 3, 'dat_nespresso_profile_s', ['标签', '标签值'], ['当月店铺购买人群', '同比月店铺购买人群', '当月行业购买人群', '同比月行业购买人群', '同比月对比行业tgi', '当月对比行业tgi', '100'], 20, '咖啡机')
    format_chart_line(prs, 16, 3, 6)
    format_chart_title(prs, 16, 3, '咖啡机购买人群画像（Nespresso官旗vs.行业）')
    format_table_data_profile(prs, fdata, 16, 4, 'dat_nespresso_profile_s', ['标签值'], ['店铺人数同比', '行业人数同比'], 20, '咖啡机')
    format_remark_text(prs, 16, 2, '策略中心')

    format_chart_data(prs, fdata, 17, 3, 'dat_nespresso_profile_d', ['标签', '标签值'], ['店铺购买人群', '店铺品牌新客', '店铺品牌老客'], 19, '咖啡机')
    format_chart_title(prs, 17, 3, '咖啡机新老客购买人群画像')
    format_remark_text(prs, 17, 2, '达摩盘', '店铺咖啡机品牌新客-当月购买店铺咖啡机差前365天购买品牌产品；\n店铺咖啡机品牌老客-当月购买店铺咖啡机交前365天购买品牌产品')

    format_chart_data_wb(prs, fdata, 18, 3, '咖啡机表现', 'B:C', 180, 11)
    format_table_data_wb(prs, fdata, 18, 4, '咖啡机表现', 'A:F', 180, 11, 1, 0)
    format_table_data_wb(prs, fdata, 18, 5, '咖啡机表现', 'I:O', 228, 7, 1, 1)
    format_table_cell_color(prs, 18, 5, list(range(2, 8)), list(range(6, 7)), '本品竞争力')
    format_remark_text(prs, 18, 2, '策略中心', '近半年咖啡机品牌新客：近180d店铺咖啡机类目购买 差 前365天店铺券品类购买\n新客来源：近180d店铺咖啡机品牌新客在前一个周期产生的购买行为；排名为占比超过20%类目按照TGI排序得到；\n*本品/竞品购买率：在近90天同时浏览本品&竞品胶囊咖啡类目后在近90天有该类目购买行为的人群在同时浏览人群中的占比')

    format_chart_data_wb(prs, fdata, 19, 8, '咖啡机表现', 'R:U', 195, 4)
    format_chart_data_wb(prs, fdata, 19, 5, '咖啡机表现', 'M:O', 244, 6)
    format_chart_data_wb(prs, fdata, 19, 4, '咖啡机表现', 'M:O', 253, 6)
    format_table_data_wb(prs, fdata, 19, 6, '咖啡机表现', 'Q:T', 200, 4, -1, 0)
    format_table_cell_data_val(prs, 19, 6, 1, 1, "{:.0f}")
    format_table_cell_data_val(prs, 19, 6, 1, 2, "{:.0f}")
    format_table_cell_data_val(prs, 19, 6, 2, 1, "{:.3%}")
    format_table_cell_data_val(prs, 19, 6, 2, 2, "{:.3%}")
    format_remark_text(prs, 19, 2, '策略中心&数据银行', '咖啡机型号升级：前12个月购买官旗2000元以下咖啡机 交 当月购买官旗2000元以上咖啡机；\n店铺咖啡机老客：前365天发生过店铺咖啡机购买行为\n*销量净变化：相较于MAT2022(2021.10.1 – 2022.9.30)，MAT2023(2022.10.1 – 2023.9.30) 购买Nespresso品牌销量的得失变化')

    format_chart_data_wb(prs, fdata, 20, 8, '咖啡机表现', 'C:E', 228, 11)
    format_chart_data_wb(prs, fdata, 20, 10, '咖啡机表现', 'Q:U', 213, 4, 2)
    format_arrow_wb(prs, fdata, 20, 16, '咖啡机表现', 'E:F', 228, 2, 1)
    format_arrow_wb(prs, fdata, 20, 17, '咖啡机表现', 'E:F', 229, 2, 1)
    format_arrow_wb(prs, fdata, 20, 18, '咖啡机表现', 'E:F', 230, 2, 1)
    format_arrow_wb(prs, fdata, 20, 19, '咖啡机表现', 'E:F', 231, 2, 1)
    format_arrow_wb(prs, fdata, 20, 20, '咖啡机表现', 'E:F', 232, 2, 1)
    format_arrow_wb(prs, fdata, 20, 21, '咖啡机表现', 'E:F', 233, 2, 1)
    format_arrow_wb(prs, fdata, 20, 22, '咖啡机表现', 'E:F', 234, 2, 1)
    format_arrow_wb(prs, fdata, 20, 23, '咖啡机表现', 'E:F', 235, 2, 1)
    format_arrow_wb(prs, fdata, 20, 24, '咖啡机表现', 'E:F', 236, 2, 1)
    format_arrow_wb(prs, fdata, 20, 25, '咖啡机表现', 'E:F', 237, 2, 1)
    format_remark_text(prs, 20, 7, '策略中心', '店铺咖啡机老客：前365天发生过店铺咖啡机购买行为；\n市场胶囊咖啡机老客：前365天发生过市场胶囊咖啡机购买行为')

    format_shape_text(prs, 22, 5, get_dt_year_month_last(), '')
    format_shape_text(prs, 22, 6, get_dt_year_month(), '')
    format_chart_data_wb_item(prs, fdata, 22, 3, '咖啡机表现', 'A:D', 276, 25)
    format_chart_data_wb_item(prs, fdata, 22, 4, '咖啡机表现', 'E:H', 276, 25)
    format_text_data_wb(prs, fdata, 22, 17, '咖啡机表现', 'B:C', 270, 2, 'cnt')
    format_text_data_wb(prs, fdata, 22, 18, '咖啡机表现', 'B:C', 271, 2, 'cnt')
    format_text_data_wb(prs, fdata, 22, 19, '咖啡机表现', 'B:C', 272, 2, 'cnt')
    format_text_data_wb(prs, fdata, 22, 20, '咖啡机表现', 'C:D', 270, 2, 'ratio')
    format_text_data_wb(prs, fdata, 22, 21, '咖啡机表现', 'C:D', 271, 2, 'ratio')
    format_text_data_wb(prs, fdata, 22, 22, '咖啡机表现', 'C:D', 272, 2, 'ratio')
    format_text_data_wb(prs, fdata, 22, 23, '咖啡机表现', 'D:E', 270, 2, 'cnt')
    format_text_data_wb(prs, fdata, 22, 24, '咖啡机表现', 'D:E', 271, 2, 'cnt')
    format_text_data_wb(prs, fdata, 22, 25, '咖啡机表现', 'D:E', 272, 2, 'cnt')
    format_text_data_wb(prs, fdata, 22, 26, '咖啡机表现', 'E:F', 270, 2, 'ratio')
    format_text_data_wb(prs, fdata, 22, 27, '咖啡机表现', 'E:F', 271, 2, 'ratio')
    format_text_data_wb(prs, fdata, 22, 28, '咖啡机表现', 'E:F', 272, 2, 'ratio')
    format_remark_text(prs, 22, 2, '生意参谋','本月进入Top100榜单的胶囊咖啡机产品中心想胶囊咖啡机及小牛妈妈品质生活馆售卖的Nespresso Inissia均有刷单嫌疑\n*2023年10-12月由于生意参谋市场大盘值无法转化为真实值，用市场Top50品牌GMV加总代替，即：占比 = 当月Top300中胶囊咖啡机交易金额 / 当月Top50品牌GMV加总')

    format_table_data_wb(prs, fdata, 23, 3, '咖啡机表现', 'H:N', 305, 11, 1, 0)
    format_table_data_wb(prs, fdata, 23, 4, '咖啡机表现', 'A:G', 305, 11, 1, 0)
    format_table_title(prs, 23, 3, 0, 0, get_dt_year_month_last(), ' 咖啡机TOP10单品')
    format_table_title(prs, 23, 4, 0, 0, get_dt_year_month(), ' 咖啡机TOP10单品')
    format_table_cell_color(prs, 23, 3, list(range(2, 12)), list(range(3, 6)), 'TOP单品')
    format_table_cell_color(prs, 23, 4, list(range(2, 12)), list(range(3, 6)), 'TOP单品')
    format_table_cell_color(prs, 23, 3, list(range(2, 12)), list(range(2, 3)), '咖啡机类型')
    format_table_cell_color(prs, 23, 4, list(range(2, 12)), list(range(2, 3)), '咖啡机类型')
    format_remark_text(prs, 23, 2, '策略中心；渠道：全渠道','本月进入Top100榜单的胶囊咖啡机产品中心想胶囊咖啡机及小牛妈妈品质生活馆售卖的Nespresso Inissia均有刷单嫌疑\n*2023年10-12月由于生意参谋市场大盘值无法转化为真实值，用市场Top50品牌GMV加总代替，即：占比 = 当月Top300中胶囊咖啡机交易金额 / 当月Top50品牌GMV加总')

    format_table_data_wb(prs, fdata, 24, 3, '咖啡机表现', 'A:I', 321, 6, 0, 0)
    format_table_data_wb(prs, fdata, 24, 4, '咖啡机表现', 'Q:X', 321, 6, 0, 0)
    format_table_data_wb(prs, fdata, 24, 19, '咖啡机表现', 'G:I', 321, 6, -1, -1)
    format_table_data_wb(prs, fdata, 24, 20, '咖啡机表现', 'J:L', 321, 6, -1, -1)
    format_shape_text(prs, 24, 5, get_dt_year_month(), '')
    format_shape_text(prs, 24, 6, get_dt_year_month_last(), '')
    format_text_data_wb(prs, fdata, 24, 12, '咖啡机表现', 'U:V', 331, 2, 'TOP5GMV')
    format_text_data_wb(prs, fdata, 24, 13, '咖啡机表现', 'E:F', 331, 2, 'TOP5GMV')
    format_remark_text(prs, 24, 2, '生意参谋')

    format_chart_data_sort(prs, fdata, 26, 4, 'dat_nespresso_shop_ct_chl', ['二级', '三级'], ['访客数_同比月', '访客数_本月', '支付转化率_同比月', '支付转化率_本月'], '广告流量', '咖啡机')
    col_no_pay=format_table_data_chl_sort(prs, fdata, 26, 3, 'dat_nespresso_shop_ct_chl', ['chl'], ['支付人数_同比', '访客数_同比', '支付转化率_同比'], '广告流量', '咖啡机')
    format_remark_text(prs, 26, 2, '生意参谋',  '生意参谋暂不支持查看{}渠道的后链路转化表现'.format('、'.join(col_no_pay)))

    format_chart_data2(prs, fdata, 27, 3, 'dat_nespresso_shop_ct_chl', ['二级'], ['访客数_同比月', '访客数_本月', '支付转化率_同比月', '支付转化率_本月'], '平台流量', '咖啡机')
    format_table_data_chl(prs, fdata, 27, 4, 'dat_nespresso_shop_ct_chl', ['chl'], ['支付人数_同比', '访客数_同比', '支付转化率_同比'], '平台流量', '咖啡机')
    format_remark_text(prs, 27, 2, '生意参谋')

    format_chart_data_chl(prs, fdata, 28, 3, 'dat_nespresso_compet_chl', ['店铺'], ['uv占比'], ['广告流量', '平台流量'], '咖啡机', ['格米莱旗舰店', 'Nespresso', 'delonghi德龙旗舰店', '柏翠旗舰店', 'barsetto电器旗舰店'])
    format_table_data_val(prs, fdata, 28, 5, 1, 1, 'dat_nespresso_compet_chl', 'uv', 1, '咖啡机', '格米莱旗舰店', '平台流量-汇总-汇总')
    format_table_data_val(prs, fdata, 28, 5, 1, 2, 'dat_nespresso_compet_chl', 'uv', 1, '咖啡机', 'Nespresso', '平台流量-汇总-汇总')
    format_table_data_val(prs, fdata, 28, 5, 1, 3, 'dat_nespresso_compet_chl', 'uv', 1, '咖啡机', 'delonghi德龙旗舰店', '平台流量-汇总-汇总')
    format_table_data_val(prs, fdata, 28, 5, 1, 4, 'dat_nespresso_compet_chl', 'uv', 1, '咖啡机', '柏翠旗舰店', '平台流量-汇总-汇总')
    format_table_data_val(prs, fdata, 28, 5, 1, 5, 'dat_nespresso_compet_chl', 'uv', 1, '咖啡机', 'barsetto电器旗舰店', '平台流量-汇总-汇总')
    format_chart_data2(prs, fdata, 28, 6, 'dat_nespresso_compet_chl', ['二级'], ['uv占比_一级'], '平台流量', '咖啡机', 'barsetto电器旗舰店')
    format_chart_data2(prs, fdata, 28, 7, 'dat_nespresso_compet_chl', ['二级'], ['uv占比_一级'], '平台流量', '咖啡机', '格米莱旗舰店')
    format_chart_data2(prs, fdata, 28, 8, 'dat_nespresso_compet_chl', ['二级'], ['uv占比_一级'], '平台流量', '咖啡机', '柏翠旗舰店')
    format_chart_data2(prs, fdata, 28, 9, 'dat_nespresso_compet_chl', ['二级'], ['uv占比_一级'], '平台流量', '咖啡机', 'delonghi德龙旗舰店')
    format_chart_data2(prs, fdata, 28, 10, 'dat_nespresso_compet_chl', ['二级'], ['uv占比_一级'], '平台流量', '咖啡机', 'Nespresso')
    format_table_data_ct_chl(prs, fdata, 28, 11, 'dat_nespresso_compet_chl', ['二级'], ['支付转化率'], '平台流量', '咖啡机', '格米莱旗舰店')
    format_table_data_ct_chl(prs, fdata, 28, 13, 'dat_nespresso_compet_chl', ['二级'], ['支付转化率'], '平台流量', '咖啡机', 'Nespresso')
    format_table_data_ct_chl(prs, fdata, 28, 15, 'dat_nespresso_compet_chl', ['二级'], ['支付转化率'], '平台流量', '咖啡机', 'delonghi德龙旗舰店')
    format_table_data_ct_chl(prs, fdata, 28, 17, 'dat_nespresso_compet_chl', ['二级'], ['支付转化率'], '平台流量', '咖啡机', '柏翠旗舰店')
    format_table_data_ct_chl(prs, fdata, 28, 19, 'dat_nespresso_compet_chl', ['二级'], ['支付转化率'], '平台流量', '咖啡机', 'barsetto电器旗舰店')
    format_table_data_val(prs, fdata, 28, 21, 1, 1, 'dat_nespresso_compet', 'uv_本月', 1, '咖啡机', '格米莱旗舰店')
    format_table_data_val(prs, fdata, 28, 21, 1, 2, 'dat_nespresso_compet', 'uv_本月', 1, '咖啡机', 'Nespresso')
    format_table_data_val(prs, fdata, 28, 21, 1, 3, 'dat_nespresso_compet', 'uv_本月', 1, '咖啡机', 'delonghi德龙旗舰店')
    format_table_data_val(prs, fdata, 28, 21, 1, 4, 'dat_nespresso_compet', 'uv_本月', 1, '咖啡机', '柏翠旗舰店')
    format_table_data_val(prs, fdata, 28, 21, 1, 5, 'dat_nespresso_compet', 'uv_本月', 1, '咖啡机', 'barsetto电器旗舰店')
    format_remark_text(prs, 28, 2, '生意参谋-品类360')

    format_table_data_ct_chl_ad(prs, fdata, 29, 3, 'dat_nespresso_compet_chl', ['二级', '三级'], ['uv占比_一级', '支付转化率'], '广告流量', '咖啡机', ['格米莱旗舰店', 'Nespresso', 'delonghi德龙旗舰店', '柏翠旗舰店', 'barsetto电器旗舰店'])
    format_table_data_val(prs, fdata, 29, 3, 3, 2, 'dat_nespresso_compet_chl', 'uv', 1, '咖啡机', '格米莱旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 29, 3, 3, 3, 'dat_nespresso_compet_chl', '支付转化率', 1, '咖啡机', '格米莱旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 29, 3, 3, 4, 'dat_nespresso_compet_chl', 'uv', 1, '咖啡机', 'Nespresso', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 29, 3, 3, 5, 'dat_nespresso_compet_chl', '支付转化率', 1, '咖啡机', 'Nespresso', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 29, 3, 3, 6, 'dat_nespresso_compet_chl', 'uv', 1, '咖啡机', 'delonghi德龙旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 29, 3, 3, 7, 'dat_nespresso_compet_chl', '支付转化率', 1, '咖啡机', 'delonghi德龙旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 29, 3, 3, 8, 'dat_nespresso_compet_chl', 'uv', 1, '咖啡机', '柏翠旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 29, 3, 3, 9, 'dat_nespresso_compet_chl', '支付转化率', 1, '咖啡机', '柏翠旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 29, 3, 3, 10, 'dat_nespresso_compet_chl', 'uv', 1, '咖啡机', 'barsetto电器旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 29, 3, 3, 11, 'dat_nespresso_compet_chl', '支付转化率', 1, '咖啡机', 'barsetto电器旗舰店', '广告流量-汇总-汇总')
    format_remark_text(prs, 29, 2, '生意参谋-竞店分析；分析维度：咖啡机类目流量')

    format_chart_data_wb(prs, fdata, 30, 3, '咖啡机表现', 'F:G', 436, 11)
    format_chart_data_wb(prs, fdata, 30, 5, '咖啡机表现', 'B:D', 436, 11)
    format_table_data_wb(prs, fdata, 30, 4, '咖啡机表现', 'G:J', 436, 11, 1, -1)
    format_remark_text(prs, 30, 2, '生意参谋')

    format_chart_data9(prs, fdata, 34, 16, 'dat_nespresso_mkt_index', ['时间范围'], ['销售金额'], 14, '咖啡')
    format_text_text(prs, 34, 8, '销售金额')
    format_text_data(prs, fdata, 34, 17, 'dat_nespresso_mkt_index', '销售金额', 1, '咖啡')
    format_text_data(prs, fdata, 34, 19, 'dat_nespresso_mkt_index', '销量', 1, '咖啡')
    format_text_data(prs, fdata, 34, 20, 'dat_nespresso_mkt_index', '购买人数', 1, '咖啡')
    format_text_data(prs, fdata, 34, 21, 'dat_nespresso_mkt_index', '购买频次', 1, '咖啡')
    format_text_data(prs, fdata, 34, 36, 'dat_nespresso_mkt_index', '单次购买量', 1, '咖啡')
    format_text_data(prs, fdata, 34, 15, 'dat_nespresso_mkt_index', '平均价格', 1, '咖啡')
    format_text_data(prs, fdata, 34, 18, 'dat_nespresso_mkt_index', '销售金额_同比', 1, '咖啡')
    format_text_data(prs, fdata, 34, 24, 'dat_nespresso_mkt_index', '销量_同比', 1, '咖啡')
    format_text_data(prs, fdata, 34, 27, 'dat_nespresso_mkt_index', '购买人数_同比', 1, '咖啡')
    format_text_data(prs, fdata, 34, 30, 'dat_nespresso_mkt_index', '购买频次_同比', 1, '咖啡')
    format_text_data(prs, fdata, 34, 38, 'dat_nespresso_mkt_index', '单次购买量_同比', 1, '咖啡')
    format_text_data(prs, fdata, 34, 33, 'dat_nespresso_mkt_index', '平均价格_同比', 1, '咖啡')
    format_text_data_color(prs, 34, 18)
    format_text_data_color(prs, 34, 24)
    format_text_data_color(prs, 34, 27)
    format_text_data_color(prs, 34, 30)
    format_text_data_color(prs, 34, 33)
    format_text_data_color(prs, 34, 38)
    format_arrow(prs, 34, 13, 34, 18)
    format_arrow(prs, 34, 23, 34, 24)
    format_arrow(prs, 34, 26, 34, 27)
    format_arrow(prs, 34, 29, 34, 30)
    format_arrow(prs, 34, 32, 34, 33)
    format_arrow(prs, 34, 35, 34, 38)
    format_remark_text(prs, 34, 2, '生意参谋；渠道：全渠道', '平均价格(指数)：商品的平均成交价格(为准确洞察，不包含0.01元的订单数据)；人均购买量(指数)：消费者平均一次购买的商品件数（=销量(指数)/购买人数(指数)/购买频次）；购买频次：每个消费者平均购买次数')

    format_text_data(prs, fdata, 35, 5, 'dat_nespresso_mkt_index', '销售金额_速溶咖啡', 1, '咖啡')
    format_text_data(prs, fdata, 35, 9, 'dat_nespresso_mkt_index', '销售金额_咖啡豆', 1, '咖啡')
    format_text_data(prs, fdata, 35, 11, 'dat_nespresso_mkt_index', '销售金额_挂耳咖啡', 1, '咖啡')
    format_text_data(prs, fdata, 35, 13, 'dat_nespresso_mkt_index', '销售金额_咖啡粉', 1, '咖啡')
    format_text_data(prs, fdata, 35, 15, 'dat_nespresso_mkt_index', '销售金额_咖啡液', 1, '咖啡')
    format_text_data(prs, fdata, 35, 17, 'dat_nespresso_mkt_index', '销售金额_胶囊咖啡', 1, '咖啡')
    format_text_data(prs, fdata, 35, 18, 'dat_nespresso_mkt_index', '销售金额_速溶咖啡_同比', 1, '咖啡')
    format_text_data(prs, fdata, 35, 21, 'dat_nespresso_mkt_index', '销售金额_咖啡豆_同比', 1, '咖啡')
    format_text_data(prs, fdata, 35, 24, 'dat_nespresso_mkt_index', '销售金额_挂耳咖啡_同比', 1, '咖啡')
    format_text_data(prs, fdata, 35, 27, 'dat_nespresso_mkt_index', '销售金额_胶囊咖啡_同比', 1, '咖啡')
    format_text_data(prs, fdata, 35, 30, 'dat_nespresso_mkt_index', '销售金额_咖啡液_同比', 1, '咖啡')
    format_text_data(prs, fdata, 35, 33, 'dat_nespresso_mkt_index', '销售金额_咖啡粉_同比', 1, '咖啡')
    format_text_data_color(prs, 35, 18)
    format_text_data_color(prs, 35, 21)
    format_text_data_color(prs, 35, 24)
    format_text_data_color(prs, 35, 27)
    format_text_data_color(prs, 35, 30)
    format_text_data_color(prs, 35, 33)

    format_arrow(prs, 35, 7, 35, 18)
    format_arrow(prs, 35, 20, 35, 21)
    format_arrow(prs, 35, 23, 35, 24)
    format_arrow(prs, 35, 26, 35, 27)
    format_arrow(prs, 35, 29, 35, 30)
    format_arrow(prs, 35, 32, 35, 33)

    format_remark_text(prs, 35, 3, '生意参谋')


    format_chart_data9(prs, fdata, 37, 16, 'dat_nespresso_mkt_index', ['时间范围'], ['销售金额'], 14, '胶囊咖啡')
    format_text_text(prs, 37, 8, '销售金额')
    format_text_data(prs, fdata, 37, 17, 'dat_nespresso_mkt_index', '销售金额', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 37, 19, 'dat_nespresso_mkt_index', '销量', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 37, 20, 'dat_nespresso_mkt_index', '购买人数', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 37, 21, 'dat_nespresso_mkt_index', '购买频次', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 37, 36, 'dat_nespresso_mkt_index', '单次购买量', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 37, 15, 'dat_nespresso_mkt_index', '平均价格', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 37, 18, 'dat_nespresso_mkt_index', '销售金额_同比', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 37, 24, 'dat_nespresso_mkt_index', '销量_同比', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 37, 27, 'dat_nespresso_mkt_index', '购买人数_同比', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 37, 30, 'dat_nespresso_mkt_index', '购买频次_同比', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 37, 38, 'dat_nespresso_mkt_index', '单次购买量_同比', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 37, 33, 'dat_nespresso_mkt_index', '平均价格_同比', 1, '胶囊咖啡')
    format_text_data_color(prs, 37, 18)
    format_text_data_color(prs, 37, 24)
    format_text_data_color(prs, 37, 27)
    format_text_data_color(prs, 37, 30)
    format_text_data_color(prs, 37, 33)
    format_text_data_color(prs, 37, 38)
    format_arrow(prs, 37, 13, 37, 18)
    format_arrow(prs, 37, 23, 37, 24)
    format_arrow(prs, 37, 26, 37, 27)
    format_arrow(prs, 37, 29, 37, 30)
    format_arrow(prs, 37, 32, 37, 33)
    format_arrow(prs, 37, 35, 37, 38)
    format_remark_text(prs, 37, 2, '生意参谋；渠道：全渠道', '平均价格(指数)：商品的平均成交价格(为准确洞察，不包含0.01元的订单数据)；人均购买量(指数)：消费者平均一次购买的商品件数（=销量(指数)/购买人数(指数)/购买频次）；购买频次：每个消费者平均购买次数')

    format_chart_data6(prs, fdata, 38, 7, 'dat_nespresso_mkt_rk_brand', ['品牌名称'], ['交易金额'], ['交易金额_万'], 10, '胶囊咖啡')
    format_chart_data3(prs, fdata, 38, 10, 'dat_nespresso_mkt_rk_brand', '销售金额', ['销售金额_Top1品牌_占比', '销售金额_Top2品牌_占比', '销售金额_Top3品牌_占比',  '销售金额_TOP4-10品牌_占比', '销售金额_TOP11-20品牌_占比', '销售金额_Top品牌_其他商家_占比'], 1, '胶囊咖啡')
    format_table_data_brand_mkt_share(prs, fdata, 38, 3, 'dat_nespresso_mkt_share', ['时间范围'], ['时间范围', '销售金额_top2品牌_占比', '销售金额_top10品牌_占比'], [1, 13], '胶囊咖啡')
    format_table_data_brand(prs, fdata, 38, 6, 'dat_nespresso_mkt_rk_brand', ['品牌名称'], ['交易金额_同比', '品牌市占率', '品牌市占率_同比', '购买人数', '购买人数_同比', '客单价', '客单价_同比', '访客人数', '访客人数_同比'], 10, '胶囊咖啡')
    #format_table_data_brand_mkt(prs, fdata, 38, 6, 'dat_nespresso_mkt_overview', ['时间范围'], ['销售金额_同比', '购买人数', '购买人数_同比', '客单价', '客单价_同比', '访客人数', '访客人数_同比'], 1, '胶囊咖啡')
    format_table_data_brand_rk(prs, fdata, 38, 13, 'dat_nespresso_mkt_rk_brand', ['品牌名称'], ['排名_变化'], 10, '胶囊咖啡')
    format_text_data(prs, fdata, 38, 9, 'dat_nespresso_mkt_rk_brand', 'TTL*', 1, '胶囊咖啡')
    format_remark_text(prs, 38, 2, '生意参谋；渠道：全渠道','市场TTL：当月市场Top50品牌GMV加总；*市占率 = 品牌GMV/ 当月Top50品牌GMV加总')

    format_chart_data_wb(prs, fdata, 39, 3, '咖啡市场情况', 'C:E', 103, 11)
    format_table_data_wb(prs, fdata, 39, 4, '咖啡市场情况', [0, 1, 2, 5, 6, 7], 103, 11, 1, -1)
    format_table_cell_color(prs, 39, 4, list(range(2, 12)), list(range(1, 4)), '品牌分销')
    format_remark_text(prs, 39, 2, '生意参谋')

    format_chart_data(prs, fdata, 40, 6, 'dat_nespresso_shop_ct_each', ['时间范围'], ['支付金额'], 13, '胶囊咖啡')
    format_text_text(prs, 40, 9, '销售金额')
    format_text_data(prs, fdata, 40, 16, 'dat_nespresso_shop_ct_each', '支付金额', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 40, 17, 'dat_nespresso_shop_ct_each', '支付买家数', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 40, 18, 'dat_nespresso_shop_ct_each', '商品访客数', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 40, 19, 'dat_nespresso_shop_ct_each', '支付转化率', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 40, 20, 'dat_nespresso_shop_ct_each', '客单价', 1, '胶囊咖啡')
    format_text_data(prs, fdata, 40, 29, 'dat_nespresso_shop_ct_each', '支付金额_同比', 1, '胶囊咖啡')
    format_text_data_color(prs, 40, 29)
    format_text_data(prs, fdata, 40, 30, 'dat_nespresso_shop_ct_each', '支付买家数_同比', 1, '胶囊咖啡')
    format_text_data_color(prs, 40, 30)
    format_text_data(prs, fdata, 40, 31, 'dat_nespresso_shop_ct_each', '商品访客数_同比', 1, '胶囊咖啡')
    format_text_data_color(prs, 40, 31)
    format_text_data(prs, fdata, 40, 32, 'dat_nespresso_shop_ct_each', '支付转化率_同比', 1, '胶囊咖啡')
    format_text_data_color(prs, 40, 32)
    format_text_data(prs, fdata, 40, 33, 'dat_nespresso_shop_ct_each', '客单价_同比', 1, '胶囊咖啡')
    format_text_data_color(prs, 40, 33)
    format_text_data(prs, fdata, 40, 34, 'dat_nespresso_mkt_index', '销售金额_同比', 1, '胶囊咖啡')
    format_text_data_color(prs, 40, 34)
    format_text_data(prs, fdata, 40, 35, 'dat_nespresso_mkt_index', '购买人数_同比', 1, '胶囊咖啡')
    format_text_data_color(prs, 40, 35)
    format_arrow(prs, 40, 15, 40, 29)
    format_arrow(prs, 40, 22, 40, 30)
    format_arrow(prs, 40, 24, 40, 31)
    format_arrow(prs, 40, 26, 40, 32)
    format_arrow(prs, 40, 28, 40, 33)
    format_remark_text(prs, 40, 2, '生意参谋')

    format_chart_title(prs, 42, 3, '胶囊咖啡购买人群画像（Nespresso官旗vs.行业）')
    format_chart_data(prs, fdata, 42, 3, 'dat_nespresso_profile_s', ['标签', '标签值'], ['当月店铺购买人群', '同比月店铺购买人群', '当月行业购买人群', '同比月行业购买人群', '同比月对比行业tgi', '当月对比行业tgi', '100'], 20, '胶囊咖啡')
    format_chart_line(prs, 42, 3, 6)
    format_table_data_profile(prs, fdata, 42, 4, 'dat_nespresso_profile_s', ['标签值'], ['店铺人数同比', '行业人数同比'], 20, '胶囊咖啡')
    format_remark_text(prs, 42, 2, '策略中心')

    format_chart_title(prs, 43, 3, '胶囊咖啡新老客购买人群画像')
    format_chart_data(prs, fdata, 43, 3, 'dat_nespresso_profile_d', ['标签', '标签值'],
                      ['店铺购买人群', '店铺品牌新客', '店铺品牌老客'], 19, '胶囊咖啡')
    format_remark_text(prs, 43, 2, '达摩盘', '店铺胶囊咖啡品牌新客-当月购买店铺胶囊咖啡差前365天购买品牌产品；\n店铺胶囊咖啡品牌老客-当月购买店铺胶囊咖啡交前365天购买品牌产品')

    format_chart_data_wb(prs, fdata, 44, 3, '咖啡市场情况', 'B:C', 223, 11)
    format_table_data_wb(prs, fdata, 44, 4, '咖啡市场情况', 'A:F', 223, 11, 1, 0)
    format_remark_text(prs, 44, 2, '策略中心', '新客来源：近180d店铺胶囊咖啡品牌新客在前一个周期产生的购买行为；排名为占比超过20%类目按照TGI排序得到')

    format_table_data_wb(prs, fdata, 45, 3, '咖啡市场情况', 'L:S', 240, 4, 1, 0)
    format_table_cell_color(prs, 45, 3, list(range(2, 5)), list(range(6, 7)), '本品竞争力')
    format_remark_text(prs, 45, 2, '策略中心', '本品/竞品购买率：在近90天同时浏览本品&竞品胶囊咖啡类目后在近90天有该类目购买行为的人群在同时浏览人群中的占比')

    format_chart_data_wb(prs, fdata, 46, 8, '咖啡市场情况', 'C:E', 284, 10)
    format_chart_data_wb(prs, fdata, 46, 11, '咖啡市场情况', 'L:O', 248, 4)
    format_chart_data_wb(prs, fdata, 46, 10, '咖啡市场情况', 'Q:T', 259, 3)
    format_arrow_wb(prs, fdata, 46, 12, '咖啡市场情况', 'E:F', 284, 2, 1)
    format_arrow_wb(prs, fdata, 46, 13, '咖啡市场情况', 'E:F', 285, 2, 1)
    format_arrow_wb(prs, fdata, 46, 14, '咖啡市场情况', 'E:F', 286, 2, 1)
    format_arrow_wb(prs, fdata, 46, 15, '咖啡市场情况', 'E:F', 287, 2, 1)
    format_arrow_wb(prs, fdata, 46, 16, '咖啡市场情况', 'E:F', 288, 2, 1)
    format_arrow_wb(prs, fdata, 46, 17, '咖啡市场情况', 'E:F', 289, 2, 1)
    format_arrow_wb(prs, fdata, 46, 18, '咖啡市场情况', 'E:F', 290, 2, 1)
    format_arrow_wb(prs, fdata, 46, 19, '咖啡市场情况', 'E:F', 291, 2, 1)
    format_arrow_wb(prs, fdata, 46, 20, '咖啡市场情况', 'E:F', 292, 2, 1)
    format_remark_text(prs, 46, 7, '策略中心&达摩盘', '胶囊咖啡品牌复购率定义：22.09-23.08购买该品牌胶囊咖啡类目老客 交 23.09购买该品牌胶囊咖啡类目老客，数据来源：策略中心\n店铺胶囊咖啡老客流向： 22.09-23.08购买Nespresso胶囊咖啡类目老客在23.09月购买市场胶囊咖啡品牌流向，数据来源：达摩盘\n老客：前365天发生过店铺胶囊咖啡/市场胶囊咖啡购买行为')

    format_table_data_wb(prs, fdata, 48, 5, '咖啡市场情况', 'G:I', 302, 11, 1, 0)
    format_table_data_wb(prs, fdata, 48, 6, '咖啡市场情况', 'A:C', 302, 11, 1, 0)
    format_table_title(prs, 48, 5, 0, 0, get_dt_year_month_last(), ' 胶囊咖啡TOP10单品')
    format_table_title(prs, 48, 6, 0, 0, get_dt_year_month(), ' 胶囊咖啡TOP10单品')
    format_text_data_wb(prs, fdata, 48, 3, '咖啡市场情况', 'I:J', 312, 2, 'Nespresso GMV占比')
    format_text_data_wb(prs, fdata, 48, 4, '咖啡市场情况', 'C:D', 312, 2, 'Nespresso GMV占比')
    format_remark_text(prs, 48, 2, '策略中心；渠道：全渠道')

    format_table_data_wb(prs, fdata, 49, 3, '咖啡市场情况', 'A:H', 319, 6, 0, 0)
    format_table_data_wb(prs, fdata, 49, 4, '咖啡市场情况', 'Q:X', 319, 6, 0, 0)
    format_table_data_wb(prs, fdata, 49, 27, '咖啡市场情况', 'G:I', 319, 6, -1, -1)
    format_table_data_wb(prs, fdata, 49, 29, '咖啡市场情况', 'J:L', 319, 6, -1, -1)
    # format_table_data_wb(prs, fdata, 49, 30, '咖啡市场情况', 'H:J', 319, 6, -1, -1)
    format_shape_text(prs, 49, 5, get_dt_year_month(), '')
    format_shape_text(prs, 49, 6, get_dt_year_month_last(), '')
    format_arrow_wb(prs, fdata, 49, 7, '咖啡市场情况', 'I:J', 319, 2, 1)
    format_arrow_wb(prs, fdata, 49, 15, '咖啡市场情况', 'I:J', 320, 2, 1)
    format_arrow_wb(prs, fdata, 49, 16, '咖啡市场情况', 'I:J', 321, 2, 1)
    format_arrow_wb(prs, fdata, 49, 17, '咖啡市场情况', 'I:J', 322, 2, 1)
    format_arrow_wb(prs, fdata, 49, 18, '咖啡市场情况', 'I:J', 323, 2, 1)
    format_remark_text(prs, 49, 2, '生意参谋')

    format_chart_data_sort(prs, fdata, 51, 4, 'dat_nespresso_shop_ct_chl', ['二级', '三级'], ['访客数_同比月', '访客数_本月', '支付转化率_同比月', '支付转化率_本月'], '广告流量', '胶囊咖啡')
    col_no_pay=format_table_data_chl_sort(prs, fdata, 51, 3, 'dat_nespresso_shop_ct_chl', ['chl'], ['支付人数_同比', '访客数_同比', '支付转化率_同比'], '广告流量', '胶囊咖啡')
    format_remark_text(prs, 51, 2, '生意参谋', '生意参谋暂不支持查看{}渠道的后链路转化表现'.format('、'.join(col_no_pay)))

    format_chart_data2(prs, fdata, 52, 3, 'dat_nespresso_shop_ct_chl', ['二级'], ['访客数_同比月', '访客数_本月', '支付转化率_同比月', '支付转化率_本月'], '平台流量', '胶囊咖啡')
    format_table_data_chl(prs, fdata, 52, 4, 'dat_nespresso_shop_ct_chl', ['chl'], ['支付人数_同比', '访客数_同比', '支付转化率_同比'], '平台流量', '胶囊咖啡')
    format_remark_text(prs, 52, 2, '生意参谋', '生意参谋暂不支持查看逛逛渠道的后链路转化表现')

    format_chart_data_chl(prs, fdata, 53, 3, 'dat_nespresso_compet_chl', ['店铺'], ['uv占比'], ['广告流量', '平台流量'], '胶囊咖啡', ['隅田川旗舰店', 'Nespresso', 'PEET’S官方旗舰店', 'Dolce Gusto官方旗舰店',     '星巴克家享咖啡旗舰店'])
    format_chart_data2(prs, fdata, 53, 6, 'dat_nespresso_compet_chl', ['二级'], ['uv占比_一级'], '平台流量', '胶囊咖啡', '星巴克家享咖啡旗舰店')
    format_chart_data2(prs, fdata, 53, 7, 'dat_nespresso_compet_chl', ['二级'], ['uv占比_一级'], '平台流量', '胶囊咖啡', '隅田川旗舰店')
    format_chart_data2(prs, fdata, 53, 8, 'dat_nespresso_compet_chl', ['二级'], ['uv占比_一级'], '平台流量', '胶囊咖啡', 'Dolce Gusto官方旗舰店')
    format_chart_data2(prs, fdata, 53, 9, 'dat_nespresso_compet_chl', ['二级'], ['uv占比_一级'], '平台流量', '胶囊咖啡', 'PEET’S官方旗舰店')
    format_chart_data2(prs, fdata, 53, 10, 'dat_nespresso_compet_chl', ['二级'], ['uv占比_一级'], '平台流量', '胶囊咖啡', 'Nespresso')
    format_table_data_ct_chl(prs, fdata, 53, 11, 'dat_nespresso_compet_chl', ['二级'], ['支付转化率'], '平台流量', '胶囊咖啡', '隅田川旗舰店')
    format_table_data_ct_chl(prs, fdata, 53, 13, 'dat_nespresso_compet_chl', ['二级'], ['支付转化率'], '平台流量', '胶囊咖啡', 'Nespresso')
    format_table_data_ct_chl(prs, fdata, 53, 15, 'dat_nespresso_compet_chl', ['二级'], ['支付转化率'], '平台流量', '胶囊咖啡', 'PEET’S官方旗舰店')
    format_table_data_ct_chl(prs, fdata, 53, 17, 'dat_nespresso_compet_chl', ['二级'], ['支付转化率'], '平台流量', '胶囊咖啡', 'Dolce Gusto官方旗舰店')
    format_table_data_ct_chl(prs, fdata, 53, 19, 'dat_nespresso_compet_chl', ['二级'], ['支付转化率'], '平台流量', '胶囊咖啡', '星巴克家享咖啡旗舰店')
    format_table_data_val(prs, fdata, 53, 5, 1, 1, 'dat_nespresso_compet_chl', 'uv', 1, '胶囊咖啡', '隅田川旗舰店', '平台流量-汇总-汇总')
    format_table_data_val(prs, fdata, 53, 5, 1, 2, 'dat_nespresso_compet_chl', 'uv', 1, '胶囊咖啡', 'Nespresso', '平台流量-汇总-汇总')
    format_table_data_val(prs, fdata, 53, 5, 1, 3, 'dat_nespresso_compet_chl', 'uv', 1, '胶囊咖啡', 'PEET’S官方旗舰店', '平台流量-汇总-汇总')
    format_table_data_val(prs, fdata, 53, 5, 1, 4, 'dat_nespresso_compet_chl', 'uv', 1, '胶囊咖啡', 'Dolce Gusto官方旗舰店', '平台流量-汇总-汇总')
    format_table_data_val(prs, fdata, 53, 5, 1, 5, 'dat_nespresso_compet_chl', 'uv', 1, '胶囊咖啡', '星巴克家享咖啡旗舰店', '平台流量-汇总-汇总')
    format_table_data_val(prs, fdata, 53, 21, 1, 1, 'dat_nespresso_compet', 'uv_本月', 1, '胶囊咖啡', '隅田川旗舰店')
    format_table_data_val(prs, fdata, 53, 21, 1, 2, 'dat_nespresso_compet', 'uv_本月', 1, '胶囊咖啡', 'Nespresso')
    format_table_data_val(prs, fdata, 53, 21, 1, 3, 'dat_nespresso_compet', 'uv_本月', 1, '胶囊咖啡', 'PEET’S官方旗舰店')
    format_table_data_val(prs, fdata, 53, 21, 1, 4, 'dat_nespresso_compet', 'uv_本月', 1, '胶囊咖啡', 'Dolce Gusto官方旗舰店')
    format_table_data_val(prs, fdata, 53, 21, 1, 5, 'dat_nespresso_compet', 'uv_本月', 1, '胶囊咖啡', '星巴克家享咖啡旗舰店')
    format_remark_text(prs, 53, 2, '生意参谋-品类360')

    format_table_data_ct_chl_ad(prs, fdata, 54, 3, 'dat_nespresso_compet_chl', ['二级', '三级'], ['uv占比_一级', '支付转化率'], '广告流量', '胶囊咖啡', ['隅田川旗舰店', 'Nespresso', 'PEET’S官方旗舰店', 'Dolce Gusto官方旗舰店', '星巴克家享咖啡旗舰店'])
    format_table_data_val(prs, fdata, 54, 3, 3, 2, 'dat_nespresso_compet_chl', 'uv', 1, '胶囊咖啡', '隅田川旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 54, 3, 3, 3, 'dat_nespresso_compet_chl', '支付转化率', 1, '胶囊咖啡', '隅田川旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 54, 3, 3, 4, 'dat_nespresso_compet_chl', 'uv', 1, '胶囊咖啡', 'Nespresso', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 54, 3, 3, 5, 'dat_nespresso_compet_chl', '支付转化率', 1, '胶囊咖啡', 'Nespresso', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 54, 3, 3, 6, 'dat_nespresso_compet_chl', 'uv', 1, '胶囊咖啡', 'PEET’S官方旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 54, 3, 3, 7, 'dat_nespresso_compet_chl', '支付转化率', 1, '胶囊咖啡', 'PEET’S官方旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 54, 3, 3, 8, 'dat_nespresso_compet_chl', 'uv', 1, '胶囊咖啡', 'Dolce Gusto官方旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 54, 3, 3, 9, 'dat_nespresso_compet_chl', '支付转化率', 1, '胶囊咖啡', 'Dolce Gusto官方旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 54, 3, 3, 10, 'dat_nespresso_compet_chl', 'uv', 1, '胶囊咖啡', '星巴克家享咖啡旗舰店', '广告流量-汇总-汇总')
    format_table_data_val(prs, fdata, 54, 3, 3, 11, 'dat_nespresso_compet_chl', '支付转化率', 1, '胶囊咖啡', '星巴克家享咖啡旗舰店', '广告流量-汇总-汇总')
    format_remark_text(prs, 54, 2, '生意参谋-竞店分析；分析维度：胶囊咖啡类目流量')

    format_chart_data_wb(prs, fdata, 55, 3, '咖啡市场情况', 'F:G', 430, 11)
    format_chart_data_wb(prs, fdata, 55, 6, '咖啡市场情况', 'B:D', 430, 11)
    format_table_data_wb(prs, fdata, 55, 4, '咖啡市场情况', 'G:J', 430, 11, 1, -1)
    format_remark_text(prs, 55, 2, '生意参谋')




def run(prs_data_load, prs_load, prs_output):
    fdata = prs_data_load
    prs = Presentation(prs_load)

    rps_format(prs, fdata)

    prs.save(prs_output)


if __name__ == '__main__':
    project_folder=os.path.dirname(os.path.dirname(os.path.realpath(__file__)))
    prs_data_load = os.path.join(project_folder,'./data/nespresso_rpt_data_202401.xlsx')
    prs_load = os.path.join(project_folder,'./data/nespresso_rpt_demo.pptx')
    prs_output = os.path.join(project_folder,'./data/nespresso_rpt_ini_202401.pptx')

    run(prs_data_load, prs_load, prs_output)
