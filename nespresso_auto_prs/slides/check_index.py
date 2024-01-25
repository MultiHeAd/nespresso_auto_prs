#用来查看demoppt里每一个模块位置
from pptx import Presentation
import os
project_folder=os.path.dirname(os.path.dirname(os.path.realpath(__file__)))
prs_load = os.path.join(project_folder,'./data/nespresso_rpt_demo.pptx')
prs = Presentation(prs_load)

for i, slide in enumerate(prs.slides):
    for j, shape in enumerate(slide.shapes):
        print(i, j, shape.name, shape.text if shape.has_text_frame else '')
